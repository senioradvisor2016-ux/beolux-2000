"""
BC2000DL — Instruktionsbog (1968 B&O-manual stil)

Bygger en .pptx som ser ut som en service-/användarmanual från 1968:
papperston, B&O-röd accent, sans-serif versal-rubriker, numrerade
röda callouts, sidnumrering nere höger.

Kör: python3 build_manual_deck.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy

# ---------- Designtokens (1968 B&O paper-style) ----------
PAPER     = RGBColor(0xF4, 0xEF, 0xE4)  # cream papper
INK       = RGBColor(0x1A, 0x1A, 0x1A)  # svart tryck
RED       = RGBColor(0xC4, 0x28, 0x20)  # B&O-röd
SOFT_GRAY = RGBColor(0x9A, 0x9A, 0x9A)  # linjer & sidnummer
LINE_GRAY = RGBColor(0xC8, 0xC0, 0xB0)  # papperslinjer

FONT_HEAD = "Helvetica"
FONT_BODY = "Helvetica"
FONT_MONO = "Courier New"

# A4-landscape proportioner via 13.33" × 7.5" (16:9 standard) men styled som papper
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Marginaler
MARGIN = Inches(0.6)

# Resurser
ROOT = Path("/Users/senioradvisor/BEOOCORD 2000 DL")
SCHEMA_DIR = ROOT / "Schema"
OUT_DIR = ROOT / "plugin" / "output"
OUT_DIR.mkdir(parents=True, exist_ok=True)

IMG_SCHEMA   = "/tmp/schema_full.jpg"
IMG_ANVISNING = "/tmp/anvisning_side1.jpg"
IMG_HERO     = "/tmp/auctionet_large_item_4447864_aaaa50174e.jpg"  # top-down vy
IMG_FRONT    = "/tmp/auctionet_large_item_4447864_9ea511348c.jpg"  # 3/4-vy
IMG_LOGO_ETC = "/tmp/auctionet_large_item_4447864_4c72216145.jpg"  # logo close-up

PCB_MIC      = str(SCHEMA_DIR / "beo2000-micpcb.jpg")
PCB_PHONO    = str(SCHEMA_DIR / "beo2000-phonopcb.jpg")
PCB_RADIO   = str(SCHEMA_DIR / "beo2000-radiopcb.jpg")
PCB_REC      = str(SCHEMA_DIR / "beo2000-recpcb.jpg")
PCB_PLAY     = str(SCHEMA_DIR / "beo2000-playpcb.jpg")
PCB_OUT      = str(SCHEMA_DIR / "beo2000-outputpcb.jpg")
PCB_PWR      = str(SCHEMA_DIR / "beo2000-pwrpcb.jpg")


# ---------- Helpers ----------
def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, x, y, w, h, text, *, size=11, bold=False, italic=False,
             color=INK, align=PP_ALIGN.LEFT, font=FONT_BODY,
             space_after=0, line_height=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True

    # första paragrafen är auto-skapad
    paragraphs_text = text.split("\n") if isinstance(text, str) else text
    for i, ptext in enumerate(paragraphs_text):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_height
        if space_after:
            p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = ptext
        f = run.font
        f.name = font
        f.size = Pt(size)
        f.bold = bold
        f.italic = italic
        f.color.rgb = color
    return tb


def add_line(slide, x1, y1, x2, y2, color=LINE_GRAY, weight=0.5):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)  # 1 = STRAIGHT
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_weight=None):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is not None:
        rect.fill.solid()
        rect.fill.fore_color.rgb = fill
    else:
        rect.fill.background()
    if line_color is not None:
        rect.line.color.rgb = line_color
        rect.line.width = Pt(line_weight or 0.5)
    else:
        rect.line.fill.background()
    rect.shadow.inherit = False
    return rect


def add_red_callout(slide, cx, cy, number, *, diameter=Inches(0.32)):
    """Röd cirkel med vitt nummer, som B&O:s callouts på Anvisning Side 1."""
    r = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        cx - diameter // 2, cy - diameter // 2, diameter, diameter
    )
    r.fill.solid()
    r.fill.fore_color.rgb = RED
    r.line.color.rgb = RED
    r.shadow.inherit = False
    tf = r.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(number)
    f = run.font
    f.name = FONT_HEAD
    f.size = Pt(10)
    f.bold = True
    f.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    return r


def add_page_chrome(slide, page_no, total_pages, section_label):
    """Topp-bar + footer som B&O-manualen."""
    set_bg(slide, PAPER)

    # Topp tunn röd linje
    add_rect(slide, MARGIN, Inches(0.45), SLIDE_W - 2 * MARGIN, Pt(1.2),
             fill=RED, line_color=None)

    # Sektionsetikett (övre vänster, ovanför linje)
    add_text(slide, MARGIN, Inches(0.18), Inches(8), Inches(0.3),
             section_label.upper(),
             size=8, bold=True, color=INK, font=FONT_HEAD)

    # Krona-symbol övre höger (representerar generisk hov-krona, ej B&O)
    add_text(slide, SLIDE_W - MARGIN - Inches(2), Inches(0.18),
             Inches(2), Inches(0.3),
             "♛ DET DANSKE KVALITETSMÆRKE",
             size=8, italic=True, align=PP_ALIGN.RIGHT,
             color=SOFT_GRAY, font=FONT_HEAD)

    # Footer linje
    add_rect(slide, MARGIN, SLIDE_H - Inches(0.5),
             SLIDE_W - 2 * MARGIN, Pt(0.5),
             fill=SOFT_GRAY, line_color=None)

    # Footer text
    add_text(slide, MARGIN, SLIDE_H - Inches(0.4),
             Inches(8), Inches(0.3),
             "BC2000DL · INSTRUKTIONSBOG · INSPIRERAD AV BeoCord 2000 DE LUXE (B&O 1968)",
             size=8, color=SOFT_GRAY, font=FONT_HEAD)

    # Sidnummer höger
    add_text(slide, SLIDE_W - MARGIN - Inches(2), SLIDE_H - Inches(0.4),
             Inches(2), Inches(0.3),
             f"SIDE {page_no} AF {total_pages}",
             size=8, color=SOFT_GRAY, align=PP_ALIGN.RIGHT, font=FONT_HEAD)


# ---------- Slides ----------
def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Helt röd bakgrund som original-manualen
    set_bg(slide, RED)

    # Krona ikonisk (Unicode)
    add_text(slide, Inches(0), Inches(1.4), SLIDE_W, Inches(1.2),
             "♛",
             size=80, color=RGBColor(0xE0, 0xC8, 0x60), align=PP_ALIGN.CENTER,
             font=FONT_HEAD)

    # B&O-stil märke: svart kvadrat med vit text
    box_w, box_h = Inches(2.4), Inches(2.4)
    bx = (SLIDE_W - box_w) // 2
    by = Inches(2.6)
    box = add_rect(slide, bx, by, box_w, box_h,
                   fill=RGBColor(0x10, 0x10, 0x10), line_color=None)

    add_text(slide, bx, by + Inches(0.55), box_w, Inches(1),
             "BC\n2000\nDL",
             size=44, bold=True, color=PAPER, align=PP_ALIGN.CENTER,
             font=FONT_HEAD, line_height=1.0)
    add_text(slide, bx, by + box_h - Inches(0.55),
             box_w, Inches(0.4),
             "Plugin-prototype",
             size=10, italic=True, color=PAPER, align=PP_ALIGN.CENTER,
             font=FONT_HEAD)

    # Stor titel under boxen
    add_text(slide, Inches(0), by + box_h + Inches(0.4),
             SLIDE_W, Inches(0.7),
             "BETJENINGSANVISNING",
             size=36, bold=True, color=PAPER, align=PP_ALIGN.CENTER,
             font=FONT_HEAD)

    # Subtitel
    add_text(slide, Inches(0), by + box_h + Inches(1.1),
             SLIDE_W, Inches(0.4),
             "PLAN-SNAPSHOT FÖR VST/AU-EMULATION",
             size=12, color=PAPER, align=PP_ALIGN.CENTER,
             font=FONT_HEAD)


def slide_forord(prs, page_no, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_chrome(slide, page_no, total, "Forord")

    # Stor rubrik
    add_text(slide, MARGIN, Inches(0.85), SLIDE_W - 2 * MARGIN, Inches(0.8),
             "FORORD",
             size=28, bold=True, color=INK, font=FONT_HEAD)

    # Sektionslinje under rubriken
    add_rect(slide, MARGIN, Inches(1.55),
             Inches(0.8), Pt(2), fill=RED, line_color=None)

    # Två-spalts text
    col_w = (SLIDE_W - 2 * MARGIN - Inches(0.6)) // 2

    left_text = (
        "Denne instruktionsbog beskriver projektet BC2000DL — en "
        "VST3/AU-plugin som emulerer Bang & Olufsens BeoCord 2000 De Luxe, "
        "lanceret 1968 i København som Jacob Jensens første store reel-to-"
        "reel-konstruktion.\n\n"
        "Dokumentet er udfærdiget som teknisk plan-snapshot for konstruktøren. "
        "Det indeholder hardware-anatomi på komponentniveau, DSP-arkitektur, "
        "UX-koncept og en eksekverbar roadmap fra Fas 0 til v1-distribution."
    )
    add_text(slide, MARGIN, Inches(2.0), col_w, Inches(4.0), left_text,
             size=11, color=INK, font=FONT_BODY, line_height=1.4,
             space_after=6)

    right_text = (
        "BeoCord 2000 De Luxe er — udover en tape-recorder — en komplet "
        "3-kanals stereomixerpult med separate optage- og gengivehoveder, "
        "germaniumtransistorer (2N2613, AC125–153, OC75, AD139, AD149), "
        "100 kHz bias og DIN-1962 EQ.\n\n"
        "BC2000DL respekterer denne arkitektur block-for-block og bygger "
        "11 genanvendelige DSP-byggesten direkte fra schemaerne — ingen "
        "black-box-emulation."
    )
    add_text(slide, MARGIN + col_w + Inches(0.6), Inches(2.0),
             col_w, Inches(4.0), right_text,
             size=11, color=INK, font=FONT_BODY, line_height=1.4,
             space_after=6)

    # Disclaimer
    add_text(slide, MARGIN, SLIDE_H - Inches(1.0),
             SLIDE_W - 2 * MARGIN, Inches(0.4),
             "BC2000DL er et selvstændigt projekt — ikke tilknyttet, godkendt eller sponsoreret af Bang & Olufsen A/S.",
             size=8, italic=True, color=SOFT_GRAY, font=FONT_BODY,
             align=PP_ALIGN.CENTER)


def slide_oversigt(prs, page_no, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_chrome(slide, page_no, total, "Oversigt")

    add_text(slide, MARGIN, Inches(0.85), SLIDE_W - 2 * MARGIN, Inches(0.8),
             "BC2000DL — OVERSIGT",
             size=28, bold=True, color=INK, font=FONT_HEAD)
    add_rect(slide, MARGIN, Inches(1.55),
             Inches(0.8), Pt(2), fill=RED, line_color=None)

    # Vänster: hero-bild
    img_w = Inches(5.8)
    img_h = Inches(4.4)
    slide.shapes.add_picture(IMG_HERO, MARGIN, Inches(2.0), img_w, img_h)

    # Caption under bilden
    add_text(slide, MARGIN, Inches(2.0) + img_h + Inches(0.05),
             img_w, Inches(0.3),
             "FIG. 1 — BeoCord 2000 De Luxe (bordmodel, teak). Referens for plugin-skin.",
             size=8, italic=True, color=SOFT_GRAY, align=PP_ALIGN.CENTER,
             font=FONT_BODY)

    # Höger: 5 stjern-punkter (som original-manualens "Nogle hovedpunkter")
    rx = MARGIN + img_w + Inches(0.5)
    rw = SLIDE_W - rx - MARGIN

    add_text(slide, rx, Inches(2.0), rw, Inches(0.3),
             "NOGLE HOVEDPUNKTER FOR BC2000DL",
             size=11, bold=True, color=INK, font=FONT_HEAD)

    points = [
        ("★", "3-kanals stereomixerpult som inputstadie — Mic + Phono + Radio mixas samtidigt, exakt som original-manualens skydepotentiometre."),
        ("★", "Separat optage- og gengivehovved-modeller med fuld DIN-1962 EQ pr. hastighed (4.75 / 9.5 / 19 cm/s)."),
        ("★", "Auto-3-speed-echo: hastighedsknappen byter ekko-tid 70 / 140 / 280 ms — feature, ej fejl."),
        ("★", "Multiplay som autentisk bounce-loop med generationsforluster, ej traditionellt tape-echo."),
        ("★", "30+ numrerade scenarier från manualen mappade direkt till plugin-presets ('1968 Manual Presets')."),
    ]
    y = Inches(2.5)
    for sym, txt in points:
        add_text(slide, rx, y, Inches(0.3), Inches(0.3), sym,
                 size=12, bold=True, color=RED, font=FONT_HEAD)
        add_text(slide, rx + Inches(0.35), y, rw - Inches(0.35), Inches(0.7),
                 txt, size=10, color=INK, font=FONT_BODY, line_height=1.3)
        y += Inches(0.7)


def slide_panel(prs, page_no, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_chrome(slide, page_no, total, "Frontpanel")

    add_text(slide, MARGIN, Inches(0.85), SLIDE_W - 2 * MARGIN, Inches(0.8),
             "FRONTPANEL — 26 KONTROLLER",
             size=28, bold=True, color=INK, font=FONT_HEAD)
    add_rect(slide, MARGIN, Inches(1.55),
             Inches(0.8), Pt(2), fill=RED, line_color=None)

    # Vänster: Anvisning Side 1-bild (full bredd)
    img_w = Inches(7.6)
    img_h = Inches(5.0)
    slide.shapes.add_picture(IMG_ANVISNING, MARGIN, Inches(2.0), img_w, img_h)

    add_text(slide, MARGIN, Inches(2.0) + img_h + Inches(0.05),
             img_w, Inches(0.3),
             "FIG. 2 — Original betjeningspanel med 26 numrerade kontroller (kilde: B&O 1968).",
             size=8, italic=True, color=SOFT_GRAY, align=PP_ALIGN.CENTER,
             font=FONT_BODY)

    # Höger: kontrolltabell
    rx = MARGIN + img_w + Inches(0.3)
    rw = SLIDE_W - rx - MARGIN

    add_text(slide, rx, Inches(2.0), rw, Inches(0.3),
             "MAPPNING TILL PLUGIN-PARAMETRAR",
             size=11, bold=True, color=INK, font=FONT_HEAD)

    rows = [
        ("1", "Speed selector (4-pos)"),
        ("2", "Geartstang (Stop/Play/FF/REW)"),
        ("4–6", "Speaker EXT / INT / Mute"),
        ("8–9", "Treble / Bass (efter playback)"),
        ("10", "Balance"),
        ("10a", "Echo level"),
        ("11", "Momentanstop"),
        ("12", "Master playback volume"),
        ("13–15", "Radio / Phono / Mic gain (3-buss)"),
        ("16–17", "VU L / R (300 ms ballistik)"),
        ("18", "Echo enable"),
        ("19–20", "Track 1 / 2 routing"),
        ("22", "Monitor (Source / Tape)"),
        ("23", "Bypass tape"),
        ("24–25", "Record arm 1 / 2"),
        ("26", "Synchroplay"),
    ]
    y = Inches(2.45)
    for n, label in rows:
        add_text(slide, rx, y, Inches(0.55), Inches(0.22), f"{n}.",
                 size=9, bold=True, color=RED, font=FONT_MONO)
        add_text(slide, rx + Inches(0.6), y, rw - Inches(0.6), Inches(0.22),
                 label, size=9, color=INK, font=FONT_BODY)
        y += Inches(0.225)


def slide_anatomi(prs, page_no, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_chrome(slide, page_no, total, "Hardware-anatomi")

    add_text(slide, MARGIN, Inches(0.85), SLIDE_W - 2 * MARGIN, Inches(0.8),
             "HARDWARE-ANATOMI — 7 MODULER",
             size=28, bold=True, color=INK, font=FONT_HEAD)
    add_rect(slide, MARGIN, Inches(1.55),
             Inches(0.8), Pt(2), fill=RED, line_color=None)

    # Schema-bild centrerad
    img_w = Inches(8.4)
    img_h = Inches(4.6)
    img_x = MARGIN
    img_y = Inches(2.0)
    slide.shapes.add_picture(IMG_SCHEMA, img_x, img_y, img_w, img_h)

    add_text(slide, img_x, img_y + img_h + Inches(0.05),
             img_w, Inches(0.3),
             "FIG. 3 — Komplet schema, BeoCord 2000 De Luxe Type 4119 (B&O 1968).",
             size=8, italic=True, color=SOFT_GRAY, align=PP_ALIGN.CENTER,
             font=FONT_BODY)

    # Höger: 7 moduler-lista
    rx = img_x + img_w + Inches(0.3)
    rw = SLIDE_W - rx - MARGIN

    add_text(slide, rx, Inches(2.0), rw, Inches(0.3),
             "ALLA 7 MODULER KARTLAGDA",
             size=11, bold=True, color=INK, font=FONT_HEAD)

    modules = [
        ("8904002", "PHONO PREAMP", "UW0029 + 2N2613 + RIAA H/L"),
        ("8904003", "RADIO PREAMP", "UW0029 + 2N2613 + flat EQ"),
        ("8904004", "MIC PREAMP", "Trafo 8012003 + UW0029 + 2N2613"),
        ("8004005", "RECORD AMP", "AC126×2 + 2N2613 + DIN-1962"),
        ("8004006", "PLAYBACK AMP", "AC126×2 + 2N2613 + DIN-1962"),
        ("8004008", "PSU + BIAS-OSC", "2× AC128 → 100 kHz idealiserad"),
        ("8004014", "POWER AMP", "AC153 → AC127+AC132 → 2× AD139"),
    ]
    y = Inches(2.4)
    for code, name, desc in modules:
        add_text(slide, rx, y, Inches(0.85), Inches(0.2), code,
                 size=9, bold=True, color=RED, font=FONT_MONO)
        add_text(slide, rx + Inches(0.9), y, rw - Inches(0.9), Inches(0.2),
                 name, size=9, bold=True, color=INK, font=FONT_HEAD)
        add_text(slide, rx, y + Inches(0.2), rw, Inches(0.22),
                 desc, size=8, color=INK, font=FONT_BODY)
        y += Inches(0.55)


def slide_input_pcbs(prs, page_no, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_chrome(slide, page_no, total, "Input-preamps")

    add_text(slide, MARGIN, Inches(0.85), SLIDE_W - 2 * MARGIN, Inches(0.8),
             "INPUT-PREAMPS — 8904002 / 03 / 04",
             size=28, bold=True, color=INK, font=FONT_HEAD)
    add_rect(slide, MARGIN, Inches(1.55),
             Inches(0.8), Pt(2), fill=RED, line_color=None)

    # 3 PCB-foton i rad
    pcb_w = Inches(3.8)
    pcb_h = Inches(2.5)
    pcb_y = Inches(2.0)
    gap = (SLIDE_W - 2 * MARGIN - 3 * pcb_w) // 2

    pcbs = [
        (PCB_MIC,   "8904004 MIC PREAMP",   "UW0029 + 2N2613 + 8012003-trafo. Stereo-asymmetri inbyggd."),
        (PCB_PHONO, "8904002 PHONO PREAMP", "UW0029 + 2N2613. HS-switch H (RIAA) / L (keramic)."),
        (PCB_RADIO, "8904003 RADIO PREAMP", "UW0029 + 2N2613. HS-switch L (3 mV) / H (100 mV)."),
    ]
    for i, (path, name, desc) in enumerate(pcbs):
        x = MARGIN + i * (pcb_w + gap)
        slide.shapes.add_picture(path, x, pcb_y, pcb_w, pcb_h)

        add_text(slide, x, pcb_y + pcb_h + Inches(0.05),
                 pcb_w, Inches(0.25),
                 f"FIG. {4+i} — {name}",
                 size=8, italic=True, color=SOFT_GRAY,
                 align=PP_ALIGN.CENTER, font=FONT_BODY)
        add_text(slide, x, pcb_y + pcb_h + Inches(0.35),
                 pcb_w, Inches(0.5), desc,
                 size=9, color=INK, align=PP_ALIGN.CENTER, font=FONT_BODY,
                 line_height=1.3)

    # Tekniska data nederst
    tech_y = pcb_y + pcb_h + Inches(1.2)
    add_text(slide, MARGIN, tech_y, SLIDE_W - 2 * MARGIN, Inches(0.3),
             "TEKNISKA DATA",
             size=11, bold=True, color=INK, font=FONT_HEAD)

    techs = [
        "Mikrofon: 50–200 Ω lo-Z (gennem trafo) — 0.5 mV/500 kΩ hi-Z (forbi trafo)",
        "Phono: 2 mV/47 kΩ magnetisk (Beogram 1000 VF) — 40 mV/4 MΩ keramisk (Beogram 1000 V)",
        "Radio: 3 mV/47 kΩ low-output — 100 mV/100 kΩ line-level",
        "Stereo-asymmetri: Vc 8 V (R) / 7.2 V (L) — autentisk 1968-tolerance, modelleras explicit",
    ]
    y = tech_y + Inches(0.3)
    for t in techs:
        add_text(slide, MARGIN + Inches(0.1), y,
                 SLIDE_W - 2 * MARGIN - Inches(0.2), Inches(0.22),
                 f"·  {t}",
                 size=9, color=INK, font=FONT_MONO)
        y += Inches(0.22)


def slide_tape_power(prs, page_no, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_chrome(slide, page_no, total, "Tape-amps & Power")

    add_text(slide, MARGIN, Inches(0.85), SLIDE_W - 2 * MARGIN, Inches(0.8),
             "TAPE-AMPS & POWER — 8004005 / 06 / 14",
             size=28, bold=True, color=INK, font=FONT_HEAD)
    add_rect(slide, MARGIN, Inches(1.55),
             Inches(0.8), Pt(2), fill=RED, line_color=None)

    # 3 PCB-foton (record/playback/output)
    pcb_w = Inches(3.8)
    pcb_h = Inches(2.5)
    pcb_y = Inches(2.0)
    gap = (SLIDE_W - 2 * MARGIN - 3 * pcb_w) // 2

    pcbs = [
        (PCB_REC,  "8004005 RECORD AMP",   "AC126×2 + 2N2613. Pre-emphasis-deck 0854661 vid input."),
        (PCB_PLAY, "8004006 PLAYBACK AMP", "AC126×2 + 2N2613. DIN-1962 EQ via switch-deck."),
        (PCB_OUT,  "8004014 POWER AMP",    "2N2613 → AC153 → AC127+AC132 → 2× AD139 klass-AB."),
    ]
    for i, (path, name, desc) in enumerate(pcbs):
        x = MARGIN + i * (pcb_w + gap)
        slide.shapes.add_picture(path, x, pcb_y, pcb_w, pcb_h)

        add_text(slide, x, pcb_y + pcb_h + Inches(0.05),
                 pcb_w, Inches(0.25),
                 f"FIG. {7+i} — {name}",
                 size=8, italic=True, color=SOFT_GRAY,
                 align=PP_ALIGN.CENTER, font=FONT_BODY)
        add_text(slide, x, pcb_y + pcb_h + Inches(0.35),
                 pcb_w, Inches(0.5), desc,
                 size=9, color=INK, align=PP_ALIGN.CENTER, font=FONT_BODY,
                 line_height=1.3)

    # Tekniska data nederst
    tech_y = pcb_y + pcb_h + Inches(1.2)
    add_text(slide, MARGIN, tech_y, SLIDE_W - 2 * MARGIN, Inches(0.3),
             "TEKNISKA DATA",
             size=11, bold=True, color=INK, font=FONT_HEAD)

    techs = [
        "Hastigheder: 19 / 9.5 / 4.75 cm/s. Bias 100 kHz @ 2.3 mA. Erase 100 kHz @ 45 mA.",
        "EQ: DIN 1962. Switch-deck 0854661 kopplar 3 olika RC-LC-nät per hastighet.",
        "S/N >55 dB @ 19 cm/s. THD <1 % @ 5 W. Output 2×8 W i 4 Ω. Line-ut 0.6 V.",
        "Power-amp: NTC 50 Ω termisk bias-kompensation. AUTOMATSIKRING soft-clip vid +14 dBu.",
    ]
    y = tech_y + Inches(0.3)
    for t in techs:
        add_text(slide, MARGIN + Inches(0.1), y,
                 SLIDE_W - 2 * MARGIN - Inches(0.2), Inches(0.22),
                 f"·  {t}",
                 size=9, color=INK, font=FONT_MONO)
        y += Inches(0.22)


def slide_dsp(prs, page_no, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_chrome(slide, page_no, total, "DSP-arkitektur")

    add_text(slide, MARGIN, Inches(0.85), SLIDE_W - 2 * MARGIN, Inches(0.8),
             "DSP-ARKITEKTUR — SIGNALKEDJA",
             size=28, bold=True, color=INK, font=FONT_HEAD)
    add_rect(slide, MARGIN, Inches(1.55),
             Inches(0.8), Pt(2), fill=RED, line_color=None)

    # ASCII-stil flowchart med rektanglar
    blocks = [
        ("MIC",   "8904004"),
        ("PHONO", "8904002"),
        ("RADIO", "8904003"),
    ]
    chain = [
        ("MIXER", "3-buss"),
        ("REC-EQ", "DIN-1962"),
        ("BIAS\n+ TAPE", "Hysteres"),
        ("WOW\n& FLUTTER", "3-speed"),
        ("PLAY-EQ", "DIN-1962"),
        ("MONITOR", "Src / Tape"),
        ("TONE", "Baxandall"),
        ("POWER", "8004014"),
    ]

    # 3 input-bussar (vänster)
    by = Inches(2.2)
    bw = Inches(1.0)
    bh = Inches(0.55)
    bx = MARGIN + Inches(0.2)
    for i, (n, c) in enumerate(blocks):
        y = by + Inches(0.2) + i * Inches(0.85)
        rect = add_rect(slide, bx, y, bw, bh,
                        fill=PAPER, line_color=INK, line_weight=0.75)
        add_text(slide, bx, y + Inches(0.05), bw, Inches(0.25), n,
                 size=10, bold=True, color=INK,
                 align=PP_ALIGN.CENTER, font=FONT_HEAD)
        add_text(slide, bx, y + Inches(0.28), bw, Inches(0.2), c,
                 size=7, color=SOFT_GRAY,
                 align=PP_ALIGN.CENTER, font=FONT_MONO)

    # Pilar från 3 bussar in i mixer
    mx = bx + bw + Inches(0.3)
    my = by + Inches(0.2 + 0.85 * 1) + bh / 2  # center på mixer
    for i in range(3):
        ay = by + Inches(0.2 + 0.85 * i) + bh / 2
        add_line(slide, bx + bw, ay, mx, my, color=INK, weight=0.75)

    # Chain (8 block) — radas ut horisontellt
    chain_x = mx
    chain_y = my - bh / 2
    chain_bw = (SLIDE_W - chain_x - MARGIN - Inches(0.2)) // 8
    chain_bh = Inches(0.7)
    for i, (n, c) in enumerate(chain):
        x = chain_x + i * chain_bw
        rect = add_rect(slide, x + Inches(0.04), chain_y, chain_bw - Inches(0.08),
                        chain_bh, fill=PAPER, line_color=INK, line_weight=0.75)
        # Highlight tape-block (bias + tape) i röd
        if "TAPE" in n:
            rect.line.color.rgb = RED
            rect.line.width = Pt(1.5)
        add_text(slide, x, chain_y + Inches(0.07),
                 chain_bw, Inches(0.35), n,
                 size=9, bold=True, color=INK,
                 align=PP_ALIGN.CENTER, font=FONT_HEAD, line_height=1.0)
        add_text(slide, x, chain_y + chain_bh - Inches(0.22),
                 chain_bw, Inches(0.2), c,
                 size=7, color=SOFT_GRAY,
                 align=PP_ALIGN.CENTER, font=FONT_MONO)
        if i > 0:
            # liten pilstreck till vänster
            add_line(slide,
                     x, chain_y + chain_bh / 2,
                     x + Inches(0.04), chain_y + chain_bh / 2,
                     color=INK, weight=0.75)

    # OUT
    out_x = chain_x + 8 * chain_bw + Inches(0.05)
    add_text(slide, out_x - Inches(0.05), chain_y, Inches(0.5), chain_bh,
             "→ OUT",
             size=10, bold=True, color=INK, font=FONT_HEAD,
             align=PP_ALIGN.LEFT)

    # 11 byggstenar i tabellform under
    add_text(slide, MARGIN, Inches(4.4), SLIDE_W - 2 * MARGIN, Inches(0.3),
             "11 ÅTERANVÄNDBARA DSP-BYGGSTENAR",
             size=11, bold=True, color=INK, font=FONT_HEAD)

    classes = [
        ("Ge2N2613Stage",            "~10 instanser",  "Asymm. soft-clip + brus, PNP germanium"),
        ("GeLowNoiseStage",          "~10 instanser",  "UW0029 (PNP) / AC126 (NPN), low-noise"),
        ("AC153Stage",               "2 instanser",    "Voltage amp i power-amp"),
        ("ComplementaryDriverPair",  "2 instanser",    "AC127 + AC132 push-pull"),
        ("AD139OutputStage",         "2 instanser",    "Klass-AB output, soft-clip"),
        ("MicTransformer8012003",    "2 instanser",    "Step-up + Jiles-Atherton-saturation"),
        ("PreEmphasisDIN1962",       "1 instans",      "3 koeff-set per hastighet"),
        ("PlaybackEQ_DIN1962",       "1 instans",      "3 koeff-set, head-bump-resonans"),
        ("BiasInjection",            "1 instans",      "Idealiserad 100 kHz sinus"),
        ("TapeSaturation",           "1 instans",      "Hysteres + bias-modulation, kärnan"),
        ("WowFlutter",               "1 instans",      "3-speed, Lagrange-interpolation"),
    ]

    cy = Inches(4.75)
    col_w = (SLIDE_W - 2 * MARGIN) // 3
    for i, (cls, count, desc) in enumerate(classes):
        col = i % 3
        row = i // 3
        x = MARGIN + col * col_w
        y = cy + row * Inches(0.5)
        add_text(slide, x, y, col_w - Inches(0.2), Inches(0.22),
                 cls, size=9, bold=True, color=INK, font=FONT_MONO)
        add_text(slide, x, y + Inches(0.2), Inches(1.4), Inches(0.2),
                 count, size=8, color=RED, font=FONT_BODY, italic=True)
        add_text(slide, x + Inches(1.4), y + Inches(0.2),
                 col_w - Inches(1.6), Inches(0.2),
                 desc, size=8, color=INK, font=FONT_BODY)


def slide_ux(prs, page_no, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_chrome(slide, page_no, total, "UX — Hybrid Hero")

    add_text(slide, MARGIN, Inches(0.85), SLIDE_W - 2 * MARGIN, Inches(0.8),
             "UX-KONCEPT — HYBRID HERO",
             size=28, bold=True, color=INK, font=FONT_HEAD)
    add_rect(slide, MARGIN, Inches(1.55),
             Inches(0.8), Pt(2), fill=RED, line_color=None)

    # Vänster: hero-foto
    img_w = Inches(6.5)
    img_h = Inches(4.5)
    slide.shapes.add_picture(IMG_FRONT, MARGIN, Inches(2.0), img_w, img_h)
    add_text(slide, MARGIN, Inches(2.0) + img_h + Inches(0.05),
             img_w, Inches(0.3),
             "FIG. 10 — Referens: bordmodel m. teak + svart anodiserad alu + brushed silver.",
             size=8, italic=True, color=SOFT_GRAY,
             align=PP_ALIGN.CENTER, font=FONT_BODY)

    # Höger: materialpalett som färgrutor
    rx = MARGIN + img_w + Inches(0.3)
    rw = SLIDE_W - rx - MARGIN

    add_text(slide, rx, Inches(2.0), rw, Inches(0.3),
             "MATERIALPALETT",
             size=11, bold=True, color=INK, font=FONT_HEAD)

    swatches = [
        (RGBColor(0x9C, 0x5F, 0x35), "Teak",           "#9C5F35"),
        (RGBColor(0x1F, 0x1F, 0x1F), "Anodiserad alu", "#1F1F1F"),
        (RGBColor(0xB5, 0xB5, 0xB5), "Brushed silver", "#B5B5B5"),
        (RGBColor(0xA0, 0xA0, 0xA0), "Push-knappar",   "#A0A0A0"),
        (RGBColor(0xF0, 0xE8, 0xD8), "VU-glas",        "#F0E8D8"),
        (RGBColor(0xC4, 0x28, 0x20), "Record red",     "#C42820"),
    ]
    sy = Inches(2.45)
    for color, name, hex_code in swatches:
        # Färgruta
        rect = add_rect(slide, rx, sy, Inches(0.5), Inches(0.4),
                        fill=color, line_color=INK, line_weight=0.5)
        add_text(slide, rx + Inches(0.6), sy + Inches(0.05),
                 Inches(2), Inches(0.2), name,
                 size=9, bold=True, color=INK, font=FONT_HEAD)
        add_text(slide, rx + Inches(0.6), sy + Inches(0.2),
                 Inches(2), Inches(0.2), hex_code,
                 size=8, color=SOFT_GRAY, font=FONT_MONO)
        sy += Inches(0.5)

    # Designprinciper
    py = Inches(5.5)
    add_text(slide, rx, py, rw, Inches(0.3),
             "JENSEN-PRINCIPER",
             size=11, bold=True, color=INK, font=FONT_HEAD)

    principles = [
        "Negativ rymd som estetik",
        "Funktionsgrupper avgränsas via spacing, ej ramar",
        "Symmetri: 5 lika faders med exakt jämn spacing",
        "Brushed metal som kontrast mot mattsvart",
    ]
    py += Inches(0.3)
    for p in principles:
        add_text(slide, rx, py, rw, Inches(0.22),
                 f"·  {p}", size=9, color=INK, font=FONT_BODY)
        py += Inches(0.22)


def slide_roadmap(prs, page_no, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_chrome(slide, page_no, total, "Roadmap")

    add_text(slide, MARGIN, Inches(0.85), SLIDE_W - 2 * MARGIN, Inches(0.8),
             "ROADMAP — 10–12 VECKOR TILL v1",
             size=28, bold=True, color=INK, font=FONT_HEAD)
    add_rect(slide, MARGIN, Inches(1.55),
             Inches(0.8), Pt(2), fill=RED, line_color=None)

    # Timeline horisontellt
    phases = [
        ("FAS 0", "KOMPONENTANALYS", "1–2 d",   "✓ KLAR — alla 7 moduler kartlagda på komponentnivå"),
        ("FAS 1", "DSP-PROTOTYP",    "1–2 v",   "Ge2N2613Stage + GeLowNoiseStage + tape-block i Python/Faust"),
        ("FAS 2", "JUCE-PORTERING",  "2–3 v",   "VST3 + AU + AAX-build, parameteröversikt, MIDI-mapping"),
        ("FAS 3", "UI / SKIN",       "2–3 v",   "Hybrid Hero med 26 kontroller + 1968 Manual Presets"),
        ("FAS 4", "VALIDERING & QA", "1 v",     "Auval clean, M1 < 5 % CPU @ 48 kHz / 128 samples"),
        ("FAS 5", "DISTRIBUTION",    "1 v",     "Code-signed installer (Mac notarized + Windows signed)"),
    ]

    pw = (SLIDE_W - 2 * MARGIN) // 6
    py = Inches(2.4)
    ph = Inches(0.9)

    # Timeline-linje
    add_rect(slide, MARGIN, py + ph + Inches(0.1),
             SLIDE_W - 2 * MARGIN, Pt(1), fill=SOFT_GRAY, line_color=None)

    for i, (label, name, dur, desc) in enumerate(phases):
        x = MARGIN + i * pw

        # Box
        box_color = RED if i == 1 else (INK if i == 0 else SOFT_GRAY)
        add_rect(slide, x + Inches(0.08), py,
                 pw - Inches(0.16), ph,
                 fill=PAPER, line_color=box_color,
                 line_weight=1.5 if i <= 1 else 0.5)

        add_text(slide, x, py + Inches(0.1), pw, Inches(0.3),
                 label,
                 size=14, bold=True, color=box_color,
                 align=PP_ALIGN.CENTER, font=FONT_HEAD)
        add_text(slide, x, py + Inches(0.4), pw, Inches(0.25),
                 name,
                 size=8, bold=True, color=INK,
                 align=PP_ALIGN.CENTER, font=FONT_HEAD)
        add_text(slide, x, py + Inches(0.65), pw, Inches(0.2),
                 dur,
                 size=8, color=SOFT_GRAY, italic=True,
                 align=PP_ALIGN.CENTER, font=FONT_MONO)

        # Marker på tidslinje
        marker = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x + pw // 2 - Inches(0.06),
            py + ph + Inches(0.04),
            Inches(0.12), Inches(0.12)
        )
        marker.fill.solid()
        marker.fill.fore_color.rgb = box_color
        marker.line.fill.background()

        # Beskrivning under
        add_text(slide, x + Inches(0.05), py + ph + Inches(0.3),
                 pw - Inches(0.1), Inches(2),
                 desc,
                 size=8, color=INK, font=FONT_BODY,
                 align=PP_ALIGN.CENTER, line_height=1.3)


def slide_servicetips(prs, page_no, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_chrome(slide, page_no, total, "Servicetips & Risker")

    add_text(slide, MARGIN, Inches(0.85), SLIDE_W - 2 * MARGIN, Inches(0.8),
             "SERVICETIPS — RISKER & BESLUT",
             size=28, bold=True, color=INK, font=FONT_HEAD)
    add_rect(slide, MARGIN, Inches(1.55),
             Inches(0.8), Pt(2), fill=RED, line_color=None)

    # Två-spalt
    col_w = (SLIDE_W - 2 * MARGIN - Inches(0.6)) // 2

    # Vänster: aktiva risker
    add_text(slide, MARGIN, Inches(2.0), col_w, Inches(0.3),
             "AKTIVA RISKER",
             size=12, bold=True, color=RED, font=FONT_HEAD)

    risks = [
        ("BIAS-HYSTERES",
         "Ikke-trivial. Förenklad anhysteretisk modell klarer 80 % av sounden men inte den karakteristiska soft-clippade övre tape-saturationen. Eskaleras kun hvis lyssningstesten kräver det."),
        ("TAPE-FORMULERING",
         "Olika band (Agfa, BASF, Scotch) ger olika ljud. Vi väljer typisk 1968-formel som default + valbar tape-typ-knapp."),
        ("ECHO/SoS/MULTIPLAY-STATE",
         "Routing-grafen er icke-trivial. Stretch-risk: +2 v om vi underskattar UI-state-management."),
        ("DISTRIBUTION-MODELL",
         "Fri / kommersiell / hybrid? Beslut innan Fas 5."),
    ]
    y = Inches(2.4)
    for name, desc in risks:
        add_text(slide, MARGIN, y, col_w, Inches(0.25),
                 f"⚠  {name}",
                 size=10, bold=True, color=INK, font=FONT_HEAD)
        add_text(slide, MARGIN + Inches(0.2), y + Inches(0.25),
                 col_w - Inches(0.2), Inches(0.7),
                 desc, size=9, color=INK, font=FONT_BODY,
                 line_height=1.35)
        y += Inches(0.95)

    # Höger: lösta beslut
    rx = MARGIN + col_w + Inches(0.6)
    add_text(slide, rx, Inches(2.0), col_w, Inches(0.3),
             "LÖSTA BESLUT",
             size=12, bold=True, color=INK, font=FONT_HEAD)

    decisions = [
        ("SCOPE v1",
         "Full funktionsmatris: tape-core + Echo + S-on-S + Multiplay + Synchroplay från release."),
        ("BRANDING",
         "Generiskt namn 'BC2000DL' / 'Danish Tape 2000'. 'Inspired by / not affiliated with B&O' i splash + installer."),
        ("RAMVERK",
         "JUCE 7+ (C++). Bästa stödet för VST3 + AU + AAX i en kodbas. Python/Faust-prototyp först."),
        ("UX-KONCEPT",
         "Hybrid Hero (Koncept C): foto-kompositioner + skarp vektor-grafik. ~1200×450 px-fönster."),
        ("UX-FÖRDJUPNING",
         "Manualens 30+ scenarier mappade direkt till plugin-presets ('1968 Manual Presets')."),
    ]
    y = Inches(2.4)
    for name, desc in decisions:
        add_text(slide, rx, y, col_w, Inches(0.25),
                 f"✓  {name}",
                 size=10, bold=True, color=INK, font=FONT_HEAD)
        add_text(slide, rx + Inches(0.2), y + Inches(0.25),
                 col_w - Inches(0.2), Inches(0.7),
                 desc, size=9, color=INK, font=FONT_BODY,
                 line_height=1.35)
        y += Inches(0.85)


# ---------- Build ----------
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Reservera blank-layout
    blank_layout = prs.slide_layouts[6]

    total = 10
    slide_cover(prs)                       # 1 cover
    slide_forord(prs,        2, total)
    slide_oversigt(prs,      3, total)
    slide_panel(prs,         4, total)
    slide_anatomi(prs,       5, total)
    slide_input_pcbs(prs,    6, total)
    slide_tape_power(prs,    7, total)
    slide_dsp(prs,           8, total)
    slide_ux(prs,            9, total)
    slide_roadmap(prs,      10, total)
    # slide_servicetips() är reservation om vi vill ha 11 sidor

    out_path = OUT_DIR / "BC2000DL_Instruktionsbog.pptx"
    prs.save(str(out_path))
    print(f"OK: {out_path}  ({out_path.stat().st_size // 1024} KB, {len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
