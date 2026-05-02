"""
BC2000DL — Retro-presentation (helt egen-ritad, vektor + matplotlib).

12 sidor i 1968 B&O-manualstil, allt ritat med pptx-shapes.
Inga foton — bara ren grafik + ett par matplotlib-PNGs för EQ-kurvor.

Kör:
  python3 diagrams/render_charts.py    # för EQ + hysteres-bilder
  python3 build_retro_deck.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN

# ---------- Designtokens (1968 B&O paper-style) ----------
PAPER     = RGBColor(0xF4, 0xEF, 0xE4)
INK       = RGBColor(0x1A, 0x1A, 0x1A)
RED       = RGBColor(0xC4, 0x28, 0x20)
SOFT_GRAY = RGBColor(0x9A, 0x9A, 0x9A)
LINE_GRAY = RGBColor(0xC8, 0xC0, 0xB0)

# Plugin UX-palett (Hybrid Hero)
TEAK_DK   = RGBColor(0x7A, 0x47, 0x22)
TEAK      = RGBColor(0x9C, 0x5F, 0x35)
TEAK_LT   = RGBColor(0xB0, 0x78, 0x45)
ALU_DK    = RGBColor(0x14, 0x14, 0x14)
ALU       = RGBColor(0x1F, 0x1F, 0x1F)
ALU_LT    = RGBColor(0x2C, 0x2C, 0x2C)
SILVER    = RGBColor(0xB5, 0xB5, 0xB5)
SILVER_LT = RGBColor(0xD0, 0xD0, 0xD0)
KNAB      = RGBColor(0xA8, 0xA8, 0xA8)
VU_GLASS  = RGBColor(0xF0, 0xE8, 0xD8)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BLACK     = RGBColor(0x08, 0x08, 0x08)

FONT_HEAD = "Helvetica"
FONT_BODY = "Helvetica"
FONT_MONO = "Courier New"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN  = Inches(0.6)

ROOT       = Path("/Users/senioradvisor/BEOOCORD 2000 DL")
OUT_DIR    = ROOT / "plugin" / "output"
DIAGRAM_DIR = ROOT / "plugin" / "diagrams" / "output"
OUT_DIR.mkdir(parents=True, exist_ok=True)

CHART_EQ        = str(DIAGRAM_DIR / "din1962_eq.png")
CHART_HYSTERESIS = str(DIAGRAM_DIR / "tape_hysteresis.png")


# ---------- Generic helpers ----------
def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_text(slide, x, y, w, h, text, *, size=11, bold=False, italic=False,
             color=INK, align=PP_ALIGN.LEFT, font=FONT_BODY,
             space_after=0, line_height=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    for attr in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, attr, Emu(0))
    tf.word_wrap = True
    paragraphs = text.split("\n") if isinstance(text, str) else text
    for i, ptext in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_height
        if space_after:
            p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = ptext
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_rect(slide, x, y, w, h, *, fill=None, line_color=None, line_weight=0.5,
             rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    rect = slide.shapes.add_shape(shape_type, x, y, w, h)
    if fill is not None:
        rect.fill.solid()
        rect.fill.fore_color.rgb = fill
    else:
        rect.fill.background()
    if line_color is not None:
        rect.line.color.rgb = line_color
        rect.line.width = Pt(line_weight)
    else:
        rect.line.fill.background()
    rect.shadow.inherit = False
    if rounded:
        # Justera corner-radius via adjustment
        try:
            rect.adjustments[0] = 0.15
        except Exception:
            pass
    return rect


def add_oval(slide, x, y, w, h, *, fill=None, line_color=None, line_weight=0.5):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color is not None:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_weight)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_line(slide, x1, y1, x2, y2, color=INK, weight=0.5):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


# ---------- Page chrome (1968 B&O-stil) ----------
def add_chrome(slide, page_no, total, section):
    set_bg(slide, PAPER)
    # Topp-linje (röd)
    add_rect(slide, MARGIN, Inches(0.45),
             SLIDE_W - 2 * MARGIN, Pt(1.2), fill=RED)
    # Sektionsetikett
    add_text(slide, MARGIN, Inches(0.18), Inches(8), Inches(0.3),
             section.upper(), size=8, bold=True,
             color=INK, font=FONT_HEAD)
    # Krona höger
    add_text(slide, SLIDE_W - MARGIN - Inches(2.5), Inches(0.18),
             Inches(2.5), Inches(0.3),
             "♛  DET DANSKE KVALITETSMÆRKE",
             size=8, italic=True, align=PP_ALIGN.RIGHT,
             color=SOFT_GRAY, font=FONT_HEAD)
    # Footer
    add_rect(slide, MARGIN, SLIDE_H - Inches(0.5),
             SLIDE_W - 2 * MARGIN, Pt(0.5), fill=SOFT_GRAY)
    add_text(slide, MARGIN, SLIDE_H - Inches(0.4), Inches(8), Inches(0.3),
             "BC2000DL · RETRO-PRESENTATION · INSPIRERAD AV BeoCord 2000 DE LUXE (B&O 1968)",
             size=8, color=SOFT_GRAY, font=FONT_HEAD)
    add_text(slide, SLIDE_W - MARGIN - Inches(2),
             SLIDE_H - Inches(0.4), Inches(2), Inches(0.3),
             f"SIDE {page_no} AF {total}",
             size=8, color=SOFT_GRAY, align=PP_ALIGN.RIGHT, font=FONT_HEAD)


def add_title(slide, title, *, sub=None):
    add_text(slide, MARGIN, Inches(0.85), SLIDE_W - 2 * MARGIN, Inches(0.7),
             title, size=26, bold=True, color=INK, font=FONT_HEAD)
    add_rect(slide, MARGIN, Inches(1.50), Inches(0.8), Pt(2), fill=RED)
    if sub:
        add_text(slide, MARGIN, Inches(1.62), SLIDE_W - 2 * MARGIN, Inches(0.3),
                 sub, size=10, italic=True, color=SOFT_GRAY, font=FONT_BODY)


# ---------- Plugin UX-mockup primitives ----------
def draw_reel(slide, cx, cy, r, *, color=SILVER, hub_color=BLACK):
    """Spole med hub (cirkulär, som tape-reel sett uppifrån)."""
    add_oval(slide, cx - r, cy - r, 2*r, 2*r,
             fill=color, line_color=ALU_DK, line_weight=0.75)
    # Hub
    hub_r = r // 4
    add_oval(slide, cx - hub_r, cy - hub_r, 2*hub_r, 2*hub_r,
             fill=hub_color, line_color=hub_color)
    # Subtle centrum-prick
    add_oval(slide, cx - hub_r // 4, cy - hub_r // 4,
             hub_r // 2, hub_r // 2, fill=SILVER_LT, line_color=SILVER_LT)


def draw_vu_meter(slide, x, y, w, h, label):
    """VU-mätare: vit/cream rektangel m. röd HIGH-zon + nål."""
    # Glas-bakgrund (rounded rect)
    add_rect(slide, x, y, w, h, fill=VU_GLASS,
             line_color=ALU, line_weight=1, rounded=True)
    # Skala-graderingar (5 markeringar)
    for i in range(6):
        gx = x + Emu(int(w * (0.1 + 0.16 * i)))
        add_line(slide, gx, y + h - Inches(0.18),
                 gx, y + h - Inches(0.10),
                 color=INK, weight=0.5)
    # Röd HIGH-zon (sista 30 %)
    high_x = x + Emu(int(w * 0.7))
    high_w = Emu(int(w * 0.25))
    add_rect(slide, high_x, y + h - Inches(0.08),
             high_w, Pt(2), fill=RED)
    # Skala-text
    add_text(slide, x, y + Inches(0.05), w, Inches(0.18),
             "NORMAL          HIGH",
             size=6, color=ALU, align=PP_ALIGN.CENTER, font=FONT_HEAD)
    # Nål (linjer som indikerar nålar; statisk position ~70 %)
    needle_x = x + Emu(int(w * 0.55))
    needle_base_x = x + Emu(int(w * 0.5))
    needle_base_y = y + h - Inches(0.05)
    add_line(slide, needle_base_x, needle_base_y,
             needle_x, y + Inches(0.18), color=INK, weight=1.2)
    # Label
    add_text(slide, x - Inches(0.25), y + Inches(0.05),
             Inches(0.25), Inches(0.2),
             label, size=8, bold=True, color=INK, font=FONT_HEAD,
             align=PP_ALIGN.CENTER)


def draw_slide_fader(slide, x, y, h, *, value=0.6, label="", track_w=Emu(64008)):
    """Vertikal slide-fader med skala 0–10."""
    # Skala-text (nummer)
    label_w = Inches(0.22)
    for i, n in enumerate(range(10, -1, -1)):
        ty = y + Emu(int(h * (i / 10)))
        add_text(slide, x - label_w - Inches(0.05),
                 ty - Inches(0.05),
                 label_w, Inches(0.16),
                 str(n), size=6, color=SOFT_GRAY, font=FONT_HEAD,
                 align=PP_ALIGN.RIGHT)
    # Spår (smal vertikal rektangel)
    track_x = x
    add_rect(slide, track_x, y, track_w, h,
             fill=ALU_LT, line_color=ALU_DK, line_weight=0.5)
    # Knopp (silver-block) vid value
    knob_w = Inches(0.32)
    knob_h = Inches(0.18)
    knob_y = y + Emu(int(h * (1 - value))) - knob_h // 2
    knob_x = track_x + track_w // 2 - knob_w // 2
    add_rect(slide, knob_x, knob_y, knob_w, knob_h,
             fill=SILVER, line_color=ALU_DK, line_weight=0.5)
    # Grip-lines på knoppen
    for i in range(3):
        ly = knob_y + Inches(0.04) + i * Inches(0.04)
        add_line(slide, knob_x + Inches(0.06), ly,
                 knob_x + knob_w - Inches(0.06), ly,
                 color=ALU_DK, weight=0.4)
    # Symbol/label längst ner
    if label:
        add_text(slide, x - Inches(0.05),
                 y + h + Inches(0.05),
                 Inches(0.35), Inches(0.18),
                 label, size=7, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font=FONT_HEAD)


def draw_push_button(slide, x, y, w, h, label="", *, pressed=False, accent=False):
    """Push-knapp: ljusgrå rektangel med tunn skugga."""
    fill = SILVER if not pressed else KNAB
    if accent and pressed:
        fill = RED
    add_rect(slide, x, y, w, h, fill=fill,
             line_color=ALU_DK, line_weight=0.5)
    # Lätt skugga längst ner (linje)
    if not pressed:
        add_line(slide, x, y + h, x + w, y + h,
                 color=ALU_DK, weight=0.6)
    # Label centrerad
    if label:
        add_text(slide, x, y + h // 2 - Inches(0.07),
                 w, Inches(0.14),
                 label, size=6, bold=True,
                 color=WHITE if pressed and accent else ALU_DK,
                 align=PP_ALIGN.CENTER, font=FONT_HEAD)


def draw_rotary(slide, cx, cy, r, *, label="", value_angle=45):
    """Rund ratt med pekare, B&O-stil."""
    # Yttre cirkel
    add_oval(slide, cx - r, cy - r, 2*r, 2*r,
             fill=KNAB, line_color=ALU_DK, line_weight=0.5)
    # Inre indikator-prick (vid pekarens position)
    import math
    rad = math.radians(value_angle - 90)  # 0 deg = top
    px = cx + int((r * 0.7) * math.cos(rad))
    py = cy + int((r * 0.7) * math.sin(rad))
    dot_r = r // 6
    add_oval(slide, px - dot_r, py - dot_r, 2*dot_r, 2*dot_r,
             fill=ALU, line_color=ALU)
    # Label under
    if label:
        add_text(slide, cx - Inches(0.6),
                 cy + r + Inches(0.05),
                 Inches(1.2), Inches(0.18),
                 label, size=7, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER, font=FONT_HEAD)


def draw_transport_lever(slide, cx, cy):
    """Y-formad transport-knapp (silver), B&O-stil."""
    # << pil vänster
    add_text(slide, cx - Inches(0.6), cy - Inches(0.1),
             Inches(0.3), Inches(0.2),
             "‹‹", size=14, bold=True, color=WHITE, font=FONT_HEAD,
             align=PP_ALIGN.CENTER)
    # >> pil höger
    add_text(slide, cx + Inches(0.3), cy - Inches(0.1),
             Inches(0.3), Inches(0.2),
             "››", size=14, bold=True, color=WHITE, font=FONT_HEAD,
             align=PP_ALIGN.CENTER)
    # Y-skaft
    base_y = cy + Inches(0.25)
    pivot_y = cy + Inches(0.05)
    add_line(slide, cx, base_y, cx, pivot_y, color=SILVER_LT, weight=4)
    # Kula på toppen
    add_oval(slide, cx - Inches(0.08), pivot_y - Inches(0.08),
             Inches(0.16), Inches(0.16),
             fill=SILVER_LT, line_color=ALU_DK, line_weight=0.4)


def draw_full_panel(slide, x, y, w, h, *, recording=False):
    """Komplett plugin-frontpanel — Hybrid Hero med teak-ram + svart panel."""

    # Trä-ram (yttre)
    teak_thickness = Inches(0.35)
    add_rect(slide, x, y, w, h, fill=TEAK,
             line_color=TEAK_DK, line_weight=1.0)
    # Inre svart panel
    panel_x = x + teak_thickness
    panel_y = y + teak_thickness
    panel_w = w - 2 * teak_thickness
    panel_h = h - 2 * teak_thickness
    add_rect(slide, panel_x, panel_y, panel_w, panel_h,
             fill=ALU, line_color=ALU_DK)

    # Topp-band: silver tape-deck-area (övre 1/3)
    deck_h = Emu(int(panel_h * 0.4))
    add_rect(slide, panel_x, panel_y, panel_w, deck_h,
             fill=SILVER, line_color=ALU_DK, line_weight=0.5)

    # Två spolar (tape reels) på topp-banded
    reel_r = Emu(int(deck_h * 0.4))
    reel_y = panel_y + deck_h // 2
    reel_left_cx = panel_x + Emu(int(panel_w * 0.22))
    reel_right_cx = panel_x + Emu(int(panel_w * 0.78))
    draw_reel(slide, reel_left_cx, reel_y, reel_r)
    draw_reel(slide, reel_right_cx, reel_y, reel_r)

    # Tape-head-cover i mitten (svart triangel/rektangel)
    head_w = Inches(1.4)
    head_h = Inches(0.7)
    head_x = panel_x + (panel_w - head_w) // 2
    head_y = panel_y + Emu(int(deck_h * 0.5))
    add_rect(slide, head_x, head_y, head_w, head_h,
             fill=ALU, line_color=ALU_DK)
    # B&O-stil krona-symbol på head-cover
    add_text(slide, head_x, head_y + Inches(0.12),
             head_w, Inches(0.3),
             "♛", size=18, color=SILVER_LT,
             align=PP_ALIGN.CENTER, font=FONT_HEAD)

    # "BEOCORD 2000 DE LUXE" → "BC2000DL"-logo under tape-deck
    logo_y = panel_y + deck_h + Inches(0.05)
    add_text(slide, panel_x, logo_y, panel_w, Inches(0.25),
             "BC2000DL  ·  DANISH TAPE 2000",
             size=11, bold=True, color=SILVER_LT,
             align=PP_ALIGN.CENTER, font=FONT_HEAD)

    # ----- Lower panel: kontroller -----
    ctrl_y = logo_y + Inches(0.4)
    ctrl_h = panel_h - (deck_h + Inches(0.4))

    # Vänster zon: VU-mätare + push-knappar
    vu_w = Inches(0.95)
    vu_h = Inches(0.55)
    vu_left_x = panel_x + Inches(0.15)
    draw_vu_meter(slide, vu_left_x, ctrl_y, vu_w, vu_h, "L")
    draw_vu_meter(slide, vu_left_x + vu_w + Inches(0.15), ctrl_y, vu_w, vu_h, "R")

    # Push-knappar (2 rader × 4 knappar)
    btn_w = Inches(0.45)
    btn_h = Inches(0.22)
    btn_gap = Inches(0.05)
    btn_y_top = ctrl_y + vu_h + Inches(0.15)
    btn_y_bot = btn_y_top + btn_h + Inches(0.08)
    btn_x_start = vu_left_x

    top_labels = ["PA", "1", "2", "SYNC"]
    bot_labels = ["MON", "L", "R", "S/S"]
    for i, lbl in enumerate(top_labels):
        bx = btn_x_start + i * (btn_w + btn_gap)
        # Förste knappen (PA) tryckt vid recording-state
        pressed = (recording and i == 1)  # spår 1 record
        accent = (i == 1)  # röd record-knapp
        draw_push_button(slide, bx, btn_y_top, btn_w, btn_h,
                         lbl, pressed=pressed, accent=accent)
    for i, lbl in enumerate(bot_labels):
        bx = btn_x_start + i * (btn_w + btn_gap)
        draw_push_button(slide, bx, btn_y_bot, btn_w, btn_h, lbl)

    # Mitten zon: 5 slide-faders
    fader_section_x = btn_x_start + 4 * (btn_w + btn_gap) + Inches(0.4)
    fader_h = Inches(1.3)
    fader_y = ctrl_y + Inches(0.15)
    fader_gap = Inches(0.55)
    fader_symbols = ["P", "G", "M", "S", "E"]  # phono/grammofon/mic/sos/echo
    fader_values = [0.6, 0.5, 0.7, 0.0, 0.4]
    for i in range(5):
        fx = fader_section_x + i * fader_gap
        draw_slide_fader(slide, fx, fader_y, fader_h,
                         value=fader_values[i],
                         label=fader_symbols[i])

    # Höger zon: transport + rotary-rattar
    right_zone_x = fader_section_x + 5 * fader_gap + Inches(0.4)
    right_zone_w = panel_w - (right_zone_x - panel_x) - Inches(0.2)

    # Transport-lever
    trans_cx = right_zone_x + right_zone_w // 4
    trans_cy = ctrl_y + Inches(0.25)
    draw_transport_lever(slide, trans_cx, trans_cy)

    # Rotary-rattar (treble + bass under)
    rotary_r = Inches(0.22)
    rot_cy_treble = ctrl_y + Inches(0.95)
    rot_cy_bass   = ctrl_y + Inches(1.55)
    rot_cx_left   = right_zone_x + Inches(0.4)
    rot_cx_right  = right_zone_x + Inches(0.95)

    draw_rotary(slide, rot_cx_left,  rot_cy_treble, rotary_r,
                label="TREBLE", value_angle=20)
    draw_rotary(slide, rot_cx_right, rot_cy_treble, rotary_r,
                label="BASS",   value_angle=-15)
    draw_rotary(slide, rot_cx_left,  rot_cy_bass, rotary_r,
                label="BAL",    value_angle=0)
    draw_rotary(slide, rot_cx_right, rot_cy_bass, rotary_r,
                label="VOL",    value_angle=60)

    # Tape-counter (3-digit) i hörnet
    counter_w = Inches(0.6)
    counter_h = Inches(0.22)
    counter_x = right_zone_x + Inches(1.5)
    counter_y = ctrl_y + Inches(0.15)
    add_rect(slide, counter_x, counter_y, counter_w, counter_h,
             fill=BLACK, line_color=ALU_DK, line_weight=0.5)
    add_text(slide, counter_x, counter_y + Inches(0.03),
             counter_w, Inches(0.18),
             "0 0 0",
             size=10, bold=True, color=RED,
             align=PP_ALIGN.CENTER, font=FONT_MONO)

    # Speed-vred (4-pos rotary)
    speed_cx = right_zone_x + Inches(1.8)
    speed_cy = ctrl_y + Inches(0.95)
    draw_rotary(slide, speed_cx, speed_cy, Inches(0.28),
                label="cm/sek", value_angle=45)

    # "BC2000DL"-textur längst ner till höger
    add_text(slide, panel_x + panel_w - Inches(0.55),
             panel_y + panel_h - Inches(0.18),
             Inches(0.5), Inches(0.13),
             "🎧",
             size=10, color=SILVER_LT, font=FONT_HEAD,
             align=PP_ALIGN.CENTER)

    # Recording-belysning bakom VU vid recording=True
    if recording:
        # Lägg en röd glow runt VU
        for vx in [vu_left_x, vu_left_x + vu_w + Inches(0.15)]:
            add_rect(slide, vx - Inches(0.05),
                     ctrl_y - Inches(0.05),
                     vu_w + Inches(0.10),
                     vu_h + Inches(0.10),
                     fill=None, line_color=RED, line_weight=2)


# ---------- Slide 1: Cover ----------
def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, RED)

    # Generic crown ornament
    add_text(slide, Inches(0), Inches(0.5), SLIDE_W, Inches(1),
             "♛", size=70, color=RGBColor(0xE0, 0xC8, 0x60),
             align=PP_ALIGN.CENTER, font=FONT_HEAD)

    # Hand-ritad mini-bandare (centrerad)
    drawing_w = Inches(7)
    drawing_h = Inches(3)
    dx = (SLIDE_W - drawing_w) // 2
    dy = Inches(1.7)

    # Träram (rektangel)
    add_rect(slide, dx, dy, drawing_w, drawing_h,
             fill=PAPER, line_color=BLACK, line_weight=2)

    # Inner svart panel
    pad = Inches(0.18)
    add_rect(slide, dx + pad, dy + pad,
             drawing_w - 2 * pad, drawing_h - 2 * pad,
             fill=BLACK, line_color=BLACK)

    # Två spolar
    reel_r = Inches(0.7)
    reel_cy = dy + Inches(1.0)
    reel_l_cx = dx + Inches(1.5)
    reel_r_cx = dx + drawing_w - Inches(1.5)
    draw_reel(slide, reel_l_cx, reel_cy, reel_r,
              color=PAPER, hub_color=BLACK)
    draw_reel(slide, reel_r_cx, reel_cy, reel_r,
              color=PAPER, hub_color=BLACK)

    # Tape-band-symboler mellan spolarna
    band_y = dy + Inches(1.9)
    band_height = Inches(0.06)
    add_rect(slide, dx + Inches(0.5), band_y,
             drawing_w - Inches(1), band_height,
             fill=PAPER, line_color=PAPER)

    # 5 mini-faders (förenklade som vita streck)
    fader_y = dy + Inches(2.2)
    for i in range(5):
        fx = dx + Inches(2.2) + i * Inches(0.4)
        add_rect(slide, fx, fader_y, Inches(0.05), Inches(0.5),
                 fill=PAPER, line_color=PAPER)

    # Stor titel
    add_text(slide, Inches(0), Inches(4.95), SLIDE_W, Inches(0.9),
             "BC2000DL",
             size=56, bold=True, color=PAPER,
             align=PP_ALIGN.CENTER, font=FONT_HEAD)

    add_text(slide, Inches(0), Inches(6.05), SLIDE_W, Inches(0.4),
             "RETRO-PRESENTATION",
             size=16, color=PAPER, align=PP_ALIGN.CENTER,
             font=FONT_HEAD)

    add_text(slide, Inches(0), Inches(6.5), SLIDE_W, Inches(0.4),
             "VST/AU PLUGIN-PROTOTYPE  ·  INSPIRERAD AV BeoCord 2000 DE LUXE",
             size=10, italic=True, color=PAPER,
             align=PP_ALIGN.CENTER, font=FONT_HEAD)


# ---------- Slide 2: Signalflöde DSP ----------
def slide_signal_flow(prs, page_no, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, page_no, total, "Signalflöde")
    add_title(slide, "DSP-SIGNALKEDJA — BLOCK FÖR BLOCK",
              sub="Komponentmodell, ingen black-box. Pilar = audio-flöde.")

    # 3 input-bussar (vänster)
    busses = [("MIC",   "8904004", Inches(2.3)),
              ("PHONO", "8904002", Inches(3.0)),
              ("RADIO", "8904003", Inches(3.7))]
    bus_w = Inches(1.0)
    bus_h = Inches(0.55)
    bus_x = MARGIN + Inches(0.1)
    for name, code, by in busses:
        add_rect(slide, bus_x, by, bus_w, bus_h,
                 fill=PAPER, line_color=INK, line_weight=1)
        add_text(slide, bus_x, by + Inches(0.06), bus_w, Inches(0.22),
                 name, size=10, bold=True, color=INK,
                 align=PP_ALIGN.CENTER, font=FONT_HEAD)
        add_text(slide, bus_x, by + Inches(0.30), bus_w, Inches(0.18),
                 code, size=7, color=SOFT_GRAY,
                 align=PP_ALIGN.CENTER, font=FONT_MONO)

    # Mixer
    mixer_x = bus_x + bus_w + Inches(0.4)
    mixer_y = Inches(2.85)
    mixer_w = Inches(1.0)
    mixer_h = Inches(0.85)
    add_rect(slide, mixer_x, mixer_y, mixer_w, mixer_h,
             fill=PAPER, line_color=INK, line_weight=1)
    add_text(slide, mixer_x, mixer_y + Inches(0.15),
             mixer_w, Inches(0.22),
             "MIXER", size=11, bold=True, color=INK,
             align=PP_ALIGN.CENTER, font=FONT_HEAD)
    add_text(slide, mixer_x, mixer_y + Inches(0.42),
             mixer_w, Inches(0.18),
             "3-buss", size=8, color=SOFT_GRAY, italic=True,
             align=PP_ALIGN.CENTER, font=FONT_MONO)

    # Linjer från bussar till mixer
    for _, _, by in busses:
        add_line(slide, bus_x + bus_w, by + bus_h // 2,
                 mixer_x, mixer_y + mixer_h // 2,
                 color=INK, weight=0.75)

    # Huvudkedja: REC-EQ → BIAS+TAPE → WOW/FLUTTER → PLAY-EQ → MONITOR → TONE → POWER
    chain = [
        ("REC-EQ",   "DIN-1962",  False),
        ("BIAS\n+TAPE",   "Hysteres",  True),    # highlight
        ("WOW\n+FLUT",    "3-speed",   False),
        ("PLAY-EQ",  "DIN-1962",  False),
        ("MONITOR",  "Src/Tape",  False),
        ("TONE",     "Baxandall", False),
        ("POWER",    "8004014",   False),
    ]
    cx_start = mixer_x + mixer_w + Inches(0.3)
    cy = mixer_y - Inches(0.07)
    cw = Inches(1.05)
    ch = Inches(0.95)
    cgap = Inches(0.18)
    for i, (n, sub, highlight) in enumerate(chain):
        x = cx_start + i * (cw + cgap)
        line_color = RED if highlight else INK
        add_rect(slide, x, cy, cw, ch,
                 fill=PAPER, line_color=line_color,
                 line_weight=1.5 if highlight else 1)
        add_text(slide, x, cy + Inches(0.15), cw, Inches(0.4),
                 n, size=10, bold=True, color=INK,
                 align=PP_ALIGN.CENTER, font=FONT_HEAD,
                 line_height=1.0)
        add_text(slide, x, cy + Inches(0.65), cw, Inches(0.18),
                 sub, size=7, color=SOFT_GRAY, italic=True,
                 align=PP_ALIGN.CENTER, font=FONT_MONO)
        # Pil mellan
        if i == 0:
            # Mixer → REC-EQ
            add_line(slide, mixer_x + mixer_w, mixer_y + mixer_h // 2,
                     x, cy + ch // 2, color=INK, weight=0.75)
        else:
            prev_x = cx_start + (i - 1) * (cw + cgap)
            add_line(slide, prev_x + cw, cy + ch // 2,
                     x, cy + ch // 2, color=INK, weight=0.75)

    # OUT
    out_x = cx_start + 7 * (cw + cgap)
    out_text_y = cy + Inches(0.35)
    add_text(slide, out_x, out_text_y, Inches(0.7), Inches(0.3),
             "▶  OUT",
             size=14, bold=True, color=RED, font=FONT_HEAD)
    add_line(slide,
             cx_start + 6 * (cw + cgap) + cw,
             cy + ch // 2,
             out_x, cy + ch // 2,
             color=INK, weight=0.75)

    # Återkopplingar (echo, S-on-S)
    fb_y = cy + ch + Inches(0.7)
    add_text(slide, MARGIN, fb_y, SLIDE_W - 2 * MARGIN, Inches(0.25),
             "ÅTERKOPPLINGAR",
             size=10, bold=True, color=INK, font=FONT_HEAD)
    fb_items = [
        ("ECHO",       "PLAY-PREAMP → delay → REC-EQ. Self-osc-tröskel modelleras."),
        ("S-on-S",     "PLAY L → REC R (kanalkorsning via 0854829)."),
        ("MULTIPLAY",  "Bounce-loop spår-till-spår med generationsförluster."),
        ("SYNCHROPLAY","Monitor-tap byts från PLAY-head till REC-head (low-latency)."),
        ("BYPASS",     "Input-mixer → tone (skip tape-block)."),
    ]
    for i, (k, v) in enumerate(fb_items):
        y = fb_y + Inches(0.3) + i * Inches(0.24)
        add_text(slide, MARGIN, y, Inches(1.5), Inches(0.22),
                 f"·  {k}",
                 size=9, bold=True, color=RED, font=FONT_HEAD)
        add_text(slide, MARGIN + Inches(1.6), y, Inches(11), Inches(0.22),
                 v, size=9, color=INK, font=FONT_BODY)


# ---------- Slide 3: DIN-1962 EQ ----------
def slide_eq(prs, page_no, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, page_no, total, "DIN-1962 EQ")
    add_title(slide, "DIN-1962 EQ — RECORD/PLAYBACK PER HASTIGHET",
              sub="Pre-emphasis vid record kompenserar tape-headets HF-roll-off. De-emphasis vid playback återställer.")

    # PNG centrerad
    img_w = Inches(11)
    img_h = Inches(4.6)
    img_x = (SLIDE_W - img_w) // 2
    img_y = Inches(2.0)
    slide.shapes.add_picture(CHART_EQ, img_x, img_y, img_w, img_h)

    # Caption
    add_text(slide, MARGIN, img_y + img_h + Inches(0.1),
             SLIDE_W - 2 * MARGIN, Inches(0.4),
             "OBSERVERA — Pre-emphasis och de-emphasis är inte exakta inverser. Mismatchen i 3–8 kHz ger den karaktäristiska \"tape glow\" som är en del av analog-tape-charmen. Pluginen renormaliserar inte bort denna.",
             size=9, italic=True, color=INK, font=FONT_BODY,
             align=PP_ALIGN.CENTER, line_height=1.4)


# ---------- Slide 4: Tape hysteresis ----------
def slide_hysteresis(prs, page_no, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, page_no, total, "Tape-modell")
    add_title(slide, "TAPE B-H-HYSTERES + 100 kHz BIAS-LINJARISERING",
              sub="Kärnblocket i pluginen — modelleras som Jiles-Atherton-approximation med dither-bias.")

    img_w = Inches(11)
    img_h = Inches(4.6)
    img_x = (SLIDE_W - img_w) // 2
    img_y = Inches(2.0)
    slide.shapes.add_picture(CHART_HYSTERESIS, img_x, img_y, img_w, img_h)

    add_text(slide, MARGIN, img_y + img_h + Inches(0.1),
             SLIDE_W - 2 * MARGIN, Inches(0.4),
             "Bias 100 kHz @ 2.3 mA (servicemanual s.2) injiceras före tape-saturationen. Den modulerar genom hysteresen så att audio-bandet effektivt återstår linjärt. Utan bias = klar saturation; med bias = soft warmth.",
             size=9, italic=True, color=INK, font=FONT_BODY,
             align=PP_ALIGN.CENTER, line_height=1.4)


# ---------- Slide 5: 3-buss-mixer-topologi ----------
def slide_mixer(prs, page_no, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, page_no, total, "3-buss mixerpult")
    add_title(slide, "DEN UNIKA 3-KANALS MIXERPULTEN",
              sub="Manualens skydepotentiometre — Mic + Phono + Radio mixas samtidigt. Inga konkurrenter har detta.")

    # Tre input-spår med preamp + fader
    track_w = Inches(2.8)
    track_h = Inches(4.0)
    track_y = Inches(2.2)
    track_gap = Inches(0.4)
    total_w = 3 * track_w + 2 * track_gap
    track_x_start = (SLIDE_W - total_w) // 2 - Inches(2)

    tracks = [
        ("MIC",   "8904004",  "Trafo 8012003 +\nUW0029 + 2N2613",
         ["Lo-Z 50–200 Ω", "Hi-Z 500 kΩ", "Stereo-asymmetri"]),
        ("PHONO", "8904002",  "UW0029 + 2N2613\nm. RIAA H/L",
         ["H: magnetisk", "L: keramisk", "+20 dB @ 50 Hz"]),
        ("RADIO", "8904003",  "UW0029 + 2N2613\nflat respons",
         ["L: 3 mV / 47 kΩ", "H: 100 mV / 100 kΩ", "Inga EQ-element"]),
    ]
    for i, (name, code, topology, specs) in enumerate(tracks):
        tx = track_x_start + i * (track_w + track_gap)

        # Track-ram
        add_rect(slide, tx, track_y, track_w, track_h,
                 fill=PAPER, line_color=INK, line_weight=0.75)

        # Header
        add_rect(slide, tx, track_y, track_w, Inches(0.4),
                 fill=INK, line_color=INK)
        add_text(slide, tx, track_y + Inches(0.07),
                 track_w, Inches(0.25),
                 name, size=14, bold=True, color=PAPER,
                 align=PP_ALIGN.CENTER, font=FONT_HEAD)

        # Module code
        add_text(slide, tx, track_y + Inches(0.45),
                 track_w, Inches(0.2),
                 code, size=9, color=RED, bold=True,
                 align=PP_ALIGN.CENTER, font=FONT_MONO)

        # Topology
        add_text(slide, tx + Inches(0.2),
                 track_y + Inches(0.75),
                 track_w - Inches(0.4), Inches(0.7),
                 topology, size=9, color=INK,
                 align=PP_ALIGN.CENTER, font=FONT_BODY)

        # Specs
        for j, sp in enumerate(specs):
            sy = track_y + Inches(1.55) + j * Inches(0.25)
            add_text(slide, tx + Inches(0.3), sy,
                     track_w - Inches(0.6), Inches(0.22),
                     f"·  {sp}", size=9, color=INK, font=FONT_BODY)

        # Slide-fader nederst (vertikal)
        fader_x = tx + track_w // 2 - Inches(0.04)
        fader_y = track_y + Inches(2.5)
        fader_h = Inches(1.2)
        add_rect(slide, fader_x, fader_y, Inches(0.08), fader_h,
                 fill=ALU, line_color=ALU_DK, line_weight=0.5)
        # Knopp
        kn_y = fader_y + Inches(0.4)
        add_rect(slide, tx + track_w // 2 - Inches(0.18),
                 kn_y, Inches(0.36), Inches(0.16),
                 fill=SILVER, line_color=ALU_DK, line_weight=0.5)
        # Skala till vänster
        for n in range(11):
            sy = fader_y + Emu(int(fader_h * (n / 10)))
            add_line(slide, fader_x - Inches(0.15), sy,
                     fader_x - Inches(0.05), sy,
                     color=SOFT_GRAY, weight=0.5)

        # Pil ner mot summation-bus
        add_text(slide, tx, track_y + track_h - Inches(0.3),
                 track_w, Inches(0.25),
                 "▼", size=16, color=RED,
                 align=PP_ALIGN.CENTER, font=FONT_HEAD)

    # Summation bus
    sum_y = track_y + track_h + Inches(0.15)
    add_rect(slide, track_x_start, sum_y,
             total_w, Inches(0.5),
             fill=INK, line_color=INK)
    add_text(slide, track_x_start, sum_y + Inches(0.13),
             total_w, Inches(0.25),
             "→  SUMMATION  →  TILL TAPE-BLOCK",
             size=11, bold=True, color=PAPER,
             align=PP_ALIGN.CENTER, font=FONT_HEAD)

    # Höger sidopanel: text om unik USP
    side_x = track_x_start + total_w + Inches(0.5)
    side_w = SLIDE_W - side_x - MARGIN

    add_text(slide, side_x, Inches(2.4), side_w, Inches(0.3),
             "VARFÖR DETTA ÄR UNIKT",
             size=11, bold=True, color=INK, font=FONT_HEAD)
    add_rect(slide, side_x, Inches(2.7), Inches(0.6), Pt(1.5), fill=RED)

    pts = [
        "UAD och Waves har bara ETT input-val per instans.",
        "BeoCord 2000 DL hade mixerpulten som standard 1968 — vi behåller det.",
        "Resultat: ETT plugin-instans = stem-mix av 3 källor med 1968 färg.",
        "Bonus: matchar manualens scenarier #1, #5, #34 direkt.",
    ]
    for i, t in enumerate(pts):
        add_text(slide, side_x, Inches(3.0) + i * Inches(0.55),
                 side_w, Inches(0.5),
                 f"·  {t}", size=10, color=INK,
                 font=FONT_BODY, line_height=1.35)


# ---------- Slide 6: 7 moduler block-diagram ----------
def slide_modules(prs, page_no, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, page_no, total, "Hardware-moduler")
    add_title(slide, "7 MODULER — KOMPONENTNIVÅ-KARTLAGDA",
              sub="Varje modul motsvarar en DSP-klass i Fas 1-prototypen.")

    # 7 boxar i 2 rader
    modules = [
        ("8904002", "PHONO PREAMP",  "UW0029 + 2N2613",       "RIAA H/L"),
        ("8904003", "RADIO PREAMP",  "UW0029 + 2N2613",       "flat EQ"),
        ("8904004", "MIC PREAMP",    "Trafo + UW0029 + 2N2613", "lo/hi-Z"),
        ("8004005", "RECORD AMP",    "AC126 + AC126 + 2N2613", "DIN-1962 pre"),
        ("8004006", "PLAYBACK AMP",  "AC126 + AC126 + 2N2613", "DIN-1962 de"),
        ("8004008", "PSU + BIAS-OSC","2× AC128 push-pull",     "100 kHz · 2.3 mA"),
        ("8004014", "POWER AMP",     "AC153 + AC127+AC132 → 2× AD139", "klass-AB"),
    ]
    bw = Inches(1.55)
    bh = Inches(1.8)
    gap_x = Inches(0.15)
    bar_y_top = Inches(2.1)
    bar_y_bot = bar_y_top + bh + Inches(0.3)
    cols_top = 4
    total_w_top = cols_top * bw + (cols_top - 1) * gap_x
    start_x_top = (SLIDE_W - total_w_top) // 2

    for i, (code, name, topo, eq) in enumerate(modules):
        if i < 4:
            x = start_x_top + i * (bw + gap_x)
            y = bar_y_top
        else:
            j = i - 4
            cols_bot = 3
            total_w_bot = cols_bot * bw + (cols_bot - 1) * gap_x
            start_x_bot = (SLIDE_W - total_w_bot) // 2
            x = start_x_bot + j * (bw + gap_x)
            y = bar_y_bot

        # Modulruta
        add_rect(slide, x, y, bw, bh,
                 fill=PAPER, line_color=INK, line_weight=0.75)
        # Header band
        add_rect(slide, x, y, bw, Inches(0.35),
                 fill=INK, line_color=INK)
        add_text(slide, x, y + Inches(0.05), bw, Inches(0.25),
                 code, size=10, bold=True, color=PAPER,
                 align=PP_ALIGN.CENTER, font=FONT_MONO)
        # Namn
        add_text(slide, x, y + Inches(0.45), bw, Inches(0.3),
                 name, size=11, bold=True, color=INK,
                 align=PP_ALIGN.CENTER, font=FONT_HEAD)
        # Röd separator
        add_rect(slide, x + Inches(0.3), y + Inches(0.78),
                 bw - Inches(0.6), Pt(1), fill=RED)
        # Topologi
        add_text(slide, x + Inches(0.1), y + Inches(0.92),
                 bw - Inches(0.2), Inches(0.6),
                 topo, size=8, color=INK, font=FONT_MONO,
                 align=PP_ALIGN.CENTER, line_height=1.3)
        # EQ/spec footer
        add_rect(slide, x, y + bh - Inches(0.32),
                 bw, Inches(0.32),
                 fill=LINE_GRAY, line_color=LINE_GRAY)
        add_text(slide, x, y + bh - Inches(0.27),
                 bw, Inches(0.22),
                 eq, size=8, italic=True, color=INK,
                 align=PP_ALIGN.CENTER, font=FONT_BODY)

    # Subtitel under
    add_text(slide, MARGIN, SLIDE_H - Inches(0.85),
             SLIDE_W - 2 * MARGIN, Inches(0.3),
             "Total kodvolym uppskattning: ~3 000–5 000 LOC C++ för Fas 2-portering till JUCE.",
             size=10, color=SOFT_GRAY, italic=True,
             align=PP_ALIGN.CENTER, font=FONT_BODY)


# ---------- Slide 7: Plugin frontpanel UX (HJÄRTAT) ----------
def slide_plugin_panel(prs, page_no, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, page_no, total, "Plugin-frontpanel")
    add_title(slide, "BC2000DL — PLUGIN-FRONTPANEL  (HYBRID HERO)",
              sub="Vektor-ritad UX-mockup. ~1200×450 px-fönster i DAW. 26 kontroller från manualens Anvisning Side 1.")

    # Centrerat panel-mockup
    panel_w = Inches(11)
    panel_h = Inches(3.8)
    panel_x = (SLIDE_W - panel_w) // 2
    panel_y = Inches(2.2)
    draw_full_panel(slide, panel_x, panel_y, panel_w, panel_h)

    # Annoteringar runt panelen
    # Vänster annotering: VU + knappar
    ann_y = panel_y + panel_h + Inches(0.15)
    add_text(slide, panel_x, ann_y, Inches(3), Inches(0.3),
             "◀  Vänsterzon: VU L/R + 2×4 push-knappar",
             size=8, italic=True, color=SOFT_GRAY, font=FONT_BODY)
    add_text(slide, panel_x + Inches(3.3), ann_y, Inches(4), Inches(0.3),
             "Mitten: 5 slide-faders (P/G/M/SoS/Echo)",
             size=8, italic=True, color=SOFT_GRAY, font=FONT_BODY,
             align=PP_ALIGN.CENTER)
    add_text(slide, panel_x + Inches(7.5), ann_y, Inches(3.5), Inches(0.3),
             "Högerzon: transport + 4 rotary + counter  ▶",
             size=8, italic=True, color=SOFT_GRAY, font=FONT_BODY,
             align=PP_ALIGN.RIGHT)

    # Materialpalett-svatscher (botten)
    swatch_y = ann_y + Inches(0.4)
    swatches = [
        (TEAK, "TEAK", "9C5F35"),
        (ALU, "ANODISERAD ALU", "1F1F1F"),
        (SILVER, "BRUSHED SILVER", "B5B5B5"),
        (KNAB, "PUSH-KNAPPAR", "A8A8A8"),
        (VU_GLASS, "VU-GLAS", "F0E8D8"),
        (RED, "RECORD RED", "C42820"),
    ]
    sw_w = Inches(0.4)
    sw_h = Inches(0.3)
    sw_total_w = len(swatches) * (sw_w + Inches(1.4))
    sw_x_start = (SLIDE_W - sw_total_w) // 2
    for i, (color, name, hex_code) in enumerate(swatches):
        sx = sw_x_start + i * (sw_w + Inches(1.4))
        add_rect(slide, sx, swatch_y, sw_w, sw_h,
                 fill=color, line_color=INK, line_weight=0.4)
        add_text(slide, sx + sw_w + Inches(0.05),
                 swatch_y, Inches(1.3), Inches(0.16),
                 name, size=7, bold=True, color=INK, font=FONT_HEAD)
        add_text(slide, sx + sw_w + Inches(0.05),
                 swatch_y + Inches(0.16),
                 Inches(1.3), Inches(0.14),
                 f"#{hex_code}", size=6, color=SOFT_GRAY, font=FONT_MONO)


# ---------- Slide 8: UX-detaljer (zoom) ----------
def slide_ux_details(prs, page_no, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, page_no, total, "UX-detaljer")
    add_title(slide, "UX-DETALJER — VU, FADER, TRANSPORT",
              sub="Tre zoom-vyer av nyckelkontroller. Hand-ritade specs för animering & ballistik.")

    # 3 detaljer i rad
    detail_w = Inches(3.5)
    detail_h = Inches(3.5)
    detail_y = Inches(2.3)
    gap = (SLIDE_W - 2 * MARGIN - 3 * detail_w) // 2

    # Detalj 1: VU-mätare (zoom 4×)
    d1_x = MARGIN
    add_rect(slide, d1_x, detail_y, detail_w, detail_h,
             fill=ALU, line_color=ALU_DK, line_weight=0.5)
    # Stort VU
    vu_big_w = Inches(2.6)
    vu_big_h = Inches(1.4)
    vu_big_x = d1_x + (detail_w - vu_big_w) // 2
    vu_big_y = detail_y + Inches(0.8)
    draw_vu_meter(slide, vu_big_x, vu_big_y, vu_big_w, vu_big_h, "L")
    # Spec under
    add_text(slide, d1_x, detail_y + Inches(2.5),
             detail_w, Inches(0.3),
             "VU-MÄTARE",
             size=12, bold=True, color=PAPER,
             align=PP_ALIGN.CENTER, font=FONT_HEAD)
    specs_vu = [
        "300 ms attack/release-ballistik",
        "Sort \"OK\"-zon, röd \"OVERLOAD\" >+3 dB",
        "Belyst med röd vid record-arm",
    ]
    for i, s in enumerate(specs_vu):
        add_text(slide, d1_x + Inches(0.2),
                 detail_y + Inches(2.85) + i * Inches(0.2),
                 detail_w - Inches(0.4), Inches(0.2),
                 f"·  {s}", size=8, color=PAPER, font=FONT_BODY)

    # Detalj 2: Slide-fader (zoom)
    d2_x = d1_x + detail_w + gap
    add_rect(slide, d2_x, detail_y, detail_w, detail_h,
             fill=ALU, line_color=ALU_DK, line_weight=0.5)
    fader_big_x = d2_x + detail_w // 2 - Inches(0.04)
    fader_big_y = detail_y + Inches(0.5)
    fader_big_h = Inches(2.0)
    # Stor fader-track
    add_rect(slide, fader_big_x, fader_big_y, Inches(0.08), fader_big_h,
             fill=ALU_LT, line_color=ALU_DK, line_weight=0.5)
    # Knopp
    kn_big_y = fader_big_y + Inches(0.5)
    add_rect(slide, d2_x + detail_w // 2 - Inches(0.4),
             kn_big_y, Inches(0.8), Inches(0.3),
             fill=SILVER, line_color=ALU_DK, line_weight=1)
    # Grip-lines
    for i in range(5):
        ly = kn_big_y + Inches(0.05) + i * Inches(0.05)
        add_line(slide, d2_x + detail_w // 2 - Inches(0.3), ly,
                 d2_x + detail_w // 2 + Inches(0.3), ly,
                 color=ALU_DK, weight=0.5)
    # Skala
    for n in range(11):
        sy = fader_big_y + Emu(int(fader_big_h * (n / 10)))
        add_line(slide, d2_x + detail_w // 2 - Inches(0.5), sy,
                 d2_x + detail_w // 2 - Inches(0.4), sy,
                 color=PAPER, weight=0.5)
        add_text(slide, d2_x + detail_w // 2 - Inches(0.78),
                 sy - Inches(0.06),
                 Inches(0.28), Inches(0.14),
                 str(10 - n), size=7, color=PAPER, font=FONT_MONO,
                 align=PP_ALIGN.RIGHT)
    # Spec under
    add_text(slide, d2_x, detail_y + Inches(2.5),
             detail_w, Inches(0.3),
             "SLIDE-FADER",
             size=12, bold=True, color=PAPER,
             align=PP_ALIGN.CENTER, font=FONT_HEAD)
    specs_f = [
        "Skala 0–10 (B&O 1968)",
        "Silver-knapp m. 5 grip-lines",
        "10 % spacing exakt jämn",
    ]
    for i, s in enumerate(specs_f):
        add_text(slide, d2_x + Inches(0.2),
                 detail_y + Inches(2.85) + i * Inches(0.2),
                 detail_w - Inches(0.4), Inches(0.2),
                 f"·  {s}", size=8, color=PAPER, font=FONT_BODY)

    # Detalj 3: Transport-lever (zoom)
    d3_x = d2_x + detail_w + gap
    add_rect(slide, d3_x, detail_y, detail_w, detail_h,
             fill=ALU, line_color=ALU_DK, line_weight=0.5)
    # ‹‹ och ›› stora
    add_text(slide, d3_x, detail_y + Inches(0.7),
             detail_w, Inches(0.5),
             "‹‹           ››",
             size=32, bold=True, color=PAPER,
             align=PP_ALIGN.CENTER, font=FONT_HEAD)
    # Stor Y-form
    cx = d3_x + detail_w // 2
    cy = detail_y + Inches(1.8)
    add_line(slide, cx, cy, cx, cy + Inches(0.6),
             color=SILVER_LT, weight=8)
    # Kula
    add_oval(slide, cx - Inches(0.18), cy - Inches(0.18),
             Inches(0.36), Inches(0.36),
             fill=SILVER_LT, line_color=ALU_DK, line_weight=0.6)
    # Spec
    add_text(slide, d3_x, detail_y + Inches(2.5),
             detail_w, Inches(0.3),
             "TRANSPORT-LEVER",
             size=12, bold=True, color=PAPER,
             align=PP_ALIGN.CENTER, font=FONT_HEAD)
    specs_t = [
        "5-state: REW · STOP · PLAY · FF · PAUSE",
        "Gravitations-baserad fysik (drag-to-pos)",
        "Aluminium-silver med skugg-djup",
    ]
    for i, s in enumerate(specs_t):
        add_text(slide, d3_x + Inches(0.2),
                 detail_y + Inches(2.85) + i * Inches(0.2),
                 detail_w - Inches(0.4), Inches(0.2),
                 f"·  {s}", size=8, color=PAPER, font=FONT_BODY)


# ---------- Slide 9: Interaktionsstates ----------
def slide_states(prs, page_no, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, page_no, total, "Interaktionsstates")
    add_title(slide, "INTERAKTIONSSTATES — IDLE → RECORDING",
              sub="Pluginen reagerar visuellt: röd belysning bakom VU vid record-arm, transport-lever ändrar position.")

    # 2 mini-paneler ovanför varandra (full bredd, mer luft)
    mini_w = Inches(11.0)
    mini_h = Inches(1.5)
    mini_y_idle = Inches(2.2)
    mini_y_rec = Inches(4.1)
    mini_x = (SLIDE_W - mini_w) // 2

    # Vänster: idle
    draw_full_panel(slide, mini_x, mini_y_idle, mini_w, mini_h,
                    recording=False)
    add_text(slide, mini_x, mini_y_idle - Inches(0.32),
             mini_w, Inches(0.3),
             "IDLE — INGEN INSPELNING",
             size=11, bold=True, color=INK,
             align=PP_ALIGN.LEFT, font=FONT_HEAD)
    add_text(slide, mini_x + Inches(3.5), mini_y_idle - Inches(0.32),
             mini_w, Inches(0.3),
             "VU-mätare i svart \"NORMAL\"-zon · Transport STOP · Spolar still",
             size=9, italic=True, color=SOFT_GRAY, font=FONT_BODY,
             align=PP_ALIGN.LEFT)

    # Höger: recording
    draw_full_panel(slide, mini_x, mini_y_rec, mini_w, mini_h,
                    recording=True)
    add_text(slide, mini_x, mini_y_rec - Inches(0.32),
             mini_w, Inches(0.3),
             "RECORDING — SPÅR 1 ARMERAD",
             size=11, bold=True, color=RED,
             align=PP_ALIGN.LEFT, font=FONT_HEAD)
    add_text(slide, mini_x + Inches(3.5), mini_y_rec - Inches(0.32),
             mini_w, Inches(0.3),
             "Röd halo bakom VU · Knap \"1\" tryckt+röd · Spolar roterar 1.0×",
             size=9, italic=True, color=SOFT_GRAY, font=FONT_BODY,
             align=PP_ALIGN.LEFT)

    # Animations-spec footer
    foot_y = mini_y_rec + mini_h + Inches(0.4)
    add_text(slide, MARGIN, foot_y, SLIDE_W - 2 * MARGIN, Inches(0.3),
             "ANIMATIONSSPEC",
             size=11, bold=True, color=INK, font=FONT_HEAD)
    add_rect(slide, MARGIN, foot_y + Inches(0.3),
             Inches(0.6), Pt(1.5), fill=RED)

    anim = [
        ("Spol-rotation",  "1× hastighet vid 19 cm/s, 0.5× vid 9.5, 0.25× vid 4.75. Reverse vid REW."),
        ("VU-ballistik",   "300 ms attack/release. Peak-hold 1.5 s med fade."),
        ("Push-tryck",     "2 px nedåt + 50 ms ease-in-out. Skugga blir djupare."),
        ("Record-glow",    "Röd halo bakom VU, 200 ms fade-in vid arm, fade-out vid disarm."),
    ]
    for i, (k, v) in enumerate(anim):
        y = foot_y + Inches(0.55) + i * Inches(0.25)
        add_text(slide, MARGIN, y, Inches(2), Inches(0.22),
                 f"·  {k}", size=9, bold=True, color=RED, font=FONT_HEAD)
        add_text(slide, MARGIN + Inches(2.1), y, Inches(11), Inches(0.22),
                 v, size=9, color=INK, font=FONT_BODY)


# ---------- Slide 10: 1968 Manual Presets meny ----------
def slide_presets(prs, page_no, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, page_no, total, "1968 Manual Presets")
    add_title(slide, "PRESET-MENY — \"1968 MANUAL PRESETS\"",
              sub="Direkt mappade från originalmanualens 30+ numrerade scenarier. Ingen annan plugin har detta.")

    # Mockup av preset-meny (centrerad, som dropdown)
    menu_w = Inches(7.5)
    menu_h = Inches(4.4)
    menu_x = (SLIDE_W - menu_w) // 2
    menu_y = Inches(2.0)

    add_rect(slide, menu_x, menu_y, menu_w, menu_h,
             fill=PAPER, line_color=INK, line_weight=1)

    # Header
    add_rect(slide, menu_x, menu_y, menu_w, Inches(0.5),
             fill=INK, line_color=INK)
    add_text(slide, menu_x, menu_y + Inches(0.13),
             menu_w, Inches(0.25),
             "PRESETS  ·  1968 MANUAL  ·  V1.0",
             size=12, bold=True, color=PAPER,
             align=PP_ALIGN.CENTER, font=FONT_HEAD)

    # Preset-rader
    presets = [
        ("MANUAL #1",  "MONO MIC",                  "Mic-bus solo, 19 cm/s, monitor=tape"),
        ("MANUAL #2",  "MONO MIC + ECHO",           "Echo enable, 70 ms slap-back"),
        ("MANUAL #4",  "MULTIPLAY 3-VOICE",         "Bounce-loop, +1 dB/gen brus"),
        ("MANUAL #5",  "STEREO MIC",                "Mic L+R, 19 cm/s hi-fi"),
        ("MANUAL #7",  "FM RADIO + ECHO",           "Radio-bus, slap-back på vokaler"),
        ("MANUAL #10", "STEREO RADIO",              "Radio-bus stereo, flat EQ"),
        ("MANUAL #15", "MONO PHONO",                "Phono-bus, RIAA H-läge"),
        ("MANUAL #21", "ELECTRIC GUITAR",           "Mic Hi-Z, 9.5 cm/s, +20 % bias"),
        ("MANUAL #22", "EL-GUITAR + SLAP",          "Som #21 + 70 ms echo"),
        ("MANUAL #34", "3-MIC ORKESTER",            "3 mics in mixerpult, full mix"),
        ("MODERN",     "DRUM-BUS GLUE",             "Tape-saturation @ +3 dB headroom"),
        ("MODERN",     "VOCAL WARMTH",              "Tape glow + 9.5 cm/s + 20 % bias"),
    ]
    row_h = Inches(0.30)
    list_y = menu_y + Inches(0.55)

    for i, (cat, name, desc) in enumerate(presets):
        y = list_y + i * row_h
        if i % 2 == 0:
            add_rect(slide, menu_x + Inches(0.05), y,
                     menu_w - Inches(0.1), row_h,
                     fill=LINE_GRAY, line_color=LINE_GRAY)
        # Highlight första preset (vald)
        if i == 0:
            add_rect(slide, menu_x + Inches(0.05), y,
                     menu_w - Inches(0.1), row_h,
                     fill=RED, line_color=RED)
            text_color = PAPER
        else:
            text_color = INK

        add_text(slide, menu_x + Inches(0.2), y + Inches(0.05),
                 Inches(1.4), Inches(0.22),
                 cat, size=8, bold=True, color=text_color, font=FONT_MONO)
        add_text(slide, menu_x + Inches(1.7), y + Inches(0.05),
                 Inches(2.5), Inches(0.22),
                 name, size=10, bold=True, color=text_color, font=FONT_HEAD)
        add_text(slide, menu_x + Inches(4.3), y + Inches(0.06),
                 Inches(3.2), Inches(0.22),
                 desc, size=8, italic=True, color=text_color, font=FONT_BODY)


# ---------- Slide 11: Roadmap ----------
def slide_roadmap(prs, page_no, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, page_no, total, "Roadmap")
    add_title(slide, "ROADMAP — 10–12 VECKOR TILL v1",
              sub="Fas 0 ✓ klar. Fas 1 påbörjad. Fas 2-5 framöver.")

    # Tidslinje horisontellt
    phases = [
        ("FAS 0", "KOMPONENT-\nANALYS", "1–2 d", "✓", "klar"),
        ("FAS 1", "DSP-\nPROTOTYP",     "1–2 v", "●", "pågår"),
        ("FAS 2", "JUCE-\nPORTERING",   "2–3 v", "○", "framöver"),
        ("FAS 3", "UI / SKIN",          "2–3 v", "○", "framöver"),
        ("FAS 4", "VALIDERING\n+ QA",   "1 v",   "○", "framöver"),
        ("FAS 5", "DISTRIBUTION",       "1 v",   "○", "framöver"),
    ]
    box_w = (SLIDE_W - 2 * MARGIN - Inches(0.4)) // 6
    box_y = Inches(2.5)
    box_h = Inches(2.0)

    # Tidslinje-axel
    axis_y = box_y + box_h + Inches(0.5)
    add_rect(slide, MARGIN, axis_y,
             SLIDE_W - 2 * MARGIN, Pt(2),
             fill=INK, line_color=INK)

    for i, (label, name, dur, marker, status) in enumerate(phases):
        x = MARGIN + i * (box_w + Inches(0.08))

        # Status-färger
        if status == "klar":
            border = INK
            border_w = 2
            fill = PAPER
        elif status == "pågår":
            border = RED
            border_w = 2.5
            fill = PAPER
        else:
            border = SOFT_GRAY
            border_w = 0.75
            fill = PAPER

        # Box
        add_rect(slide, x + Inches(0.05), box_y,
                 box_w - Inches(0.1), box_h,
                 fill=fill, line_color=border, line_weight=border_w)

        # Label header
        add_text(slide, x, box_y + Inches(0.15),
                 box_w, Inches(0.4),
                 label, size=18, bold=True,
                 color=border if status != "klar" else INK,
                 align=PP_ALIGN.CENTER, font=FONT_HEAD)

        # Röd separator
        add_rect(slide, x + box_w // 2 - Inches(0.3),
                 box_y + Inches(0.62),
                 Inches(0.6), Pt(1), fill=RED)

        # Namn
        add_text(slide, x + Inches(0.1),
                 box_y + Inches(0.78),
                 box_w - Inches(0.2), Inches(0.6),
                 name, size=10, bold=True, color=INK,
                 align=PP_ALIGN.CENTER, font=FONT_HEAD,
                 line_height=1.1)

        # Duration
        add_text(slide, x, box_y + box_h - Inches(0.4),
                 box_w, Inches(0.25),
                 dur, size=9, color=SOFT_GRAY, italic=True,
                 align=PP_ALIGN.CENTER, font=FONT_MONO)

        # Marker på axel
        marker_size = Inches(0.22)
        marker_x = x + box_w // 2 - marker_size // 2
        marker_y = axis_y - marker_size // 2 + Pt(1)
        marker_color = INK if status == "klar" else (RED if status == "pågår" else SOFT_GRAY)
        add_oval(slide, marker_x, marker_y, marker_size, marker_size,
                 fill=marker_color, line_color=marker_color)

        # Status-text under axel
        add_text(slide, x, axis_y + Inches(0.2),
                 box_w, Inches(0.22),
                 status.upper(), size=9, bold=True,
                 color=marker_color, font=FONT_HEAD,
                 align=PP_ALIGN.CENTER)


# ---------- Slide 12: Risker + Disclaimer ----------
def slide_risks(prs, page_no, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, page_no, total, "Risker · Beslut · Slut")
    add_title(slide, "ÖPPNA FRÅGOR — RISKER & BESLUT")

    # Två-spalt
    col_w = (SLIDE_W - 2 * MARGIN - Inches(0.6)) // 2

    # Vänster: aktiva risker
    add_text(slide, MARGIN, Inches(2.1), col_w, Inches(0.3),
             "AKTIVA RISKER",
             size=14, bold=True, color=RED, font=FONT_HEAD)
    add_rect(slide, MARGIN, Inches(2.4), Inches(0.6), Pt(2), fill=RED)

    risks = [
        ("BIAS-HYSTERES",
         "Förenklad modell klarar 80 % av sounden. Eskaleras bara om lyssningstesten kräver."),
        ("TAPE-FORMULERING",
         "Vi väljer typisk 1968-formel som default + valbar tape-typ-knapp."),
        ("ECHO-STATE",
         "Routing-grafen icke-trivial. Stretch-risk: +2 v om vi underskattar UI-state."),
        ("DISTRIBUTION",
         "Fri / kommersiell / hybrid? Beslut innan Fas 5."),
    ]
    y = Inches(2.6)
    for name, desc in risks:
        add_text(slide, MARGIN, y, col_w, Inches(0.25),
                 f"⚠  {name}",
                 size=11, bold=True, color=INK, font=FONT_HEAD)
        add_text(slide, MARGIN + Inches(0.25), y + Inches(0.27),
                 col_w - Inches(0.25), Inches(0.7),
                 desc, size=9, color=INK, font=FONT_BODY,
                 line_height=1.35)
        y += Inches(0.85)

    # Höger: lösta beslut
    rx = MARGIN + col_w + Inches(0.6)
    add_text(slide, rx, Inches(2.1), col_w, Inches(0.3),
             "LÖSTA BESLUT",
             size=14, bold=True, color=INK, font=FONT_HEAD)
    add_rect(slide, rx, Inches(2.4), Inches(0.6), Pt(2), fill=INK)

    decisions = [
        ("SCOPE v1",
         "Full funktionsmatris från release: tape + Echo + S-on-S + Multiplay + Synchroplay."),
        ("BRANDING",
         "BC2000DL / Danish Tape 2000. Inga B&O-namn i UI eller marknadsföring."),
        ("RAMVERK",
         "JUCE 7+ (C++). VST3 + AU + AAX i en kodbas. Python/Faust-prototyp först."),
        ("UX",
         "Hybrid Hero (Koncept C): foto-komp + skarp vektor. ~1200×450 px-fönster."),
    ]
    y = Inches(2.6)
    for name, desc in decisions:
        add_text(slide, rx, y, col_w, Inches(0.25),
                 f"✓  {name}",
                 size=11, bold=True, color=INK, font=FONT_HEAD)
        add_text(slide, rx + Inches(0.25), y + Inches(0.27),
                 col_w - Inches(0.25), Inches(0.7),
                 desc, size=9, color=INK, font=FONT_BODY,
                 line_height=1.35)
        y += Inches(0.85)

    # Disclaimer footer
    disc_y = SLIDE_H - Inches(1.1)
    add_rect(slide, MARGIN, disc_y, SLIDE_W - 2 * MARGIN, Pt(0.5),
             fill=SOFT_GRAY)
    add_text(slide, MARGIN, disc_y + Inches(0.1),
             SLIDE_W - 2 * MARGIN, Inches(0.3),
             "BC2000DL är ett självständigt projekt — inte tilknyttet, godkendt eller sponsoreret af Bang & Olufsen A/S.",
             size=9, italic=True, color=SOFT_GRAY,
             align=PP_ALIGN.CENTER, font=FONT_BODY)


# ---------- Build ----------
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    total = 12
    slide_cover(prs)
    slide_signal_flow(prs,    2, total)
    slide_eq(prs,             3, total)
    slide_hysteresis(prs,     4, total)
    slide_mixer(prs,          5, total)
    slide_modules(prs,        6, total)
    slide_plugin_panel(prs,   7, total)
    slide_ux_details(prs,     8, total)
    slide_states(prs,         9, total)
    slide_presets(prs,       10, total)
    slide_roadmap(prs,       11, total)
    slide_risks(prs,         12, total)

    out_path = OUT_DIR / "BC2000DL_Retro_Presentation.pptx"
    prs.save(str(out_path))
    print(f"OK: {out_path}  ({out_path.stat().st_size // 1024} KB, {len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
